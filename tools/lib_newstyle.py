#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Shared design system for the 2026 NCHU MCP Workshop slide deck.

Visual language: violet primary + pastel cards + heavy headings + minimal
top metadata bar. Build scripts under tools/ should import from this module
rather than redeclaring constants/helpers.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# Export everything (including underscore-prefixed helpers)
__all__ = [
    # palette
    "VIOLET", "VIOLET_DEEP", "ORANGE", "TEAL", "TEAL_DEEP", "PINK",
    "PINK_DEEP", "BLUE", "BLUE_DEEP",
    "VIOLET_PASTEL", "ORANGE_PASTEL", "TEAL_PASTEL", "PINK_PASTEL",
    "BLUE_PASTEL", "SLATE_PASTEL", "YELLOW_PASTEL",
    "BG_WHITE", "MIDNIGHT", "INK", "INK_SOFT", "MUTED", "ITALIC_GRAY", "HAIRLINE",
    "CODE_BG", "CODE_FG", "CODE_STRING", "CODE_KEYWORD", "CODE_COMMENT", "CODE_ORANGE",
    "FONT_TITLE", "FONT_BODY", "FONT_CODE",
    "SLIDE_W", "SLIDE_H", "S",
    # primitives + helpers
    "_I", "_P",
    "_rect", "_rounded", "_oval", "_text", "_multi", "_blank_slide",
    "metadata_bar", "slide_title", "slide_subtitle", "page_number",
    "pastel_card", "circle_number", "code_block", "callout_box",
    "make_presentation", "pastel_for",
    # enums for build scripts
    "RGBColor", "MSO_SHAPE", "PP_ALIGN", "MSO_ANCHOR",
    "Inches", "Pt", "Presentation",
]


# ── Palette ─────────────────────────────────────────────────────────
VIOLET       = RGBColor(0x7B, 0x5C, 0xF5)
VIOLET_DEEP  = RGBColor(0x5B, 0x3E, 0xD9)
ORANGE       = RGBColor(0xFB, 0x70, 0x48)
TEAL         = RGBColor(0x0E, 0xA5, 0x94)
TEAL_DEEP    = RGBColor(0x0F, 0x76, 0x6E)
PINK         = RGBColor(0xEF, 0x44, 0x44)
PINK_DEEP    = RGBColor(0xB9, 0x1C, 0x1C)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)
BLUE_DEEP    = RGBColor(0x1D, 0x4E, 0xD8)

# Pastel card fills
VIOLET_PASTEL = RGBColor(0xEE, 0xE8, 0xFF)
ORANGE_PASTEL = RGBColor(0xFF, 0xE7, 0xDB)
TEAL_PASTEL   = RGBColor(0xD4, 0xF7, 0xF0)
PINK_PASTEL   = RGBColor(0xFE, 0xE2, 0xE2)
BLUE_PASTEL   = RGBColor(0xDB, 0xEA, 0xFE)
SLATE_PASTEL  = RGBColor(0xF1, 0xF5, 0xF9)
YELLOW_PASTEL = RGBColor(0xFE, 0xF3, 0xC7)

# Neutrals
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MIDNIGHT   = RGBColor(0x0F, 0x14, 0x29)
INK        = RGBColor(0x1F, 0x29, 0x37)
INK_SOFT   = RGBColor(0x37, 0x41, 0x51)
MUTED      = RGBColor(0x9C, 0xA3, 0xAF)
ITALIC_GRAY = RGBColor(0x6B, 0x72, 0x80)
HAIRLINE   = RGBColor(0xE5, 0xE7, 0xEB)

# Code
CODE_BG       = RGBColor(0x1A, 0x21, 0x38)
CODE_FG       = RGBColor(0xE5, 0xE7, 0xEB)
CODE_STRING   = RGBColor(0xA8, 0xE1, 0xB4)
CODE_KEYWORD  = RGBColor(0xC8, 0xB6, 0xFF)
CODE_COMMENT  = RGBColor(0x9C, 0xA3, 0xAF)
CODE_ORANGE   = RGBColor(0xFB, 0x92, 0x4D)

# ── Fonts ──────────────────────────────────────────────────────────
FONT_TITLE = "Arial Black"
FONT_BODY  = "Calibri"
FONT_CODE  = "Consolas"

# Slide dim — pptx 10 × 5.625 inches; logical 13.333 × 7.5 scaled by S
SLIDE_W, SLIDE_H = 10.0, 5.625
S = 0.75


def _I(v): return Inches(v * S)
def _P(v): return Pt(v * S)


# ── Low-level primitives ───────────────────────────────────────────
def _rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _I(x), _I(y), _I(w), _I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = _P(line_w)
    sh.shadow.inherit = False
    return sh


def _rounded(slide, x, y, w, h, color, line_color=None, line_w=None, adj=0.10):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                _I(x), _I(y), _I(w), _I(h))
    sh.adjustments[0] = adj
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = _P(line_w)
    sh.shadow.inherit = False
    return sh


def _oval(slide, x, y, r, color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                _I(x - r), _I(y - r), _I(r * 2), _I(r * 2))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    sh.shadow.inherit = False
    return sh


def _text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18,
          color=INK, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = _I(0.04)
    tf.margin_top = tf.margin_bottom = _I(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = _P(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _multi(slide, x, y, w, h, paragraphs, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = _I(0.04)
    tf.margin_top = tf.margin_bottom = _I(0.02)
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if "space_after" in spec:
            p.space_after = _P(spec["space_after"])
        r = p.add_run()
        r.text = spec["text"]
        r.font.name = spec.get("font", FONT_BODY)
        r.font.size = _P(spec.get("size", 18))
        r.font.color.rgb = spec.get("color", INK)
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
    return tb


def _blank_slide(prs, bg=BG_WHITE):
    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    s = prs.slides.add_slide(layout)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    sp_tree = s.shapes._spTree
    sp_tree.remove(rect._element); sp_tree.insert(2, rect._element)
    return s


# ── Mid-level patterns ─────────────────────────────────────────────
def metadata_bar(slide, num, label, accent=VIOLET):
    """Top-left small colored dot + 'NUM · LABEL' uppercase letter-spaced."""
    _oval(slide, 0.85, 0.55, 0.20, accent)
    _text(slide, 1.2, 0.30, 12, 0.55,
          f"{num}   {label}",
          font=FONT_BODY, size=15, color=MUTED, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)


def slide_title(slide, text, *, y=0.95, size=42, color=INK, accent=None):
    c = accent if accent else color
    _text(slide, 0.85, y, 12, 0.95, text,
          font=FONT_TITLE, size=size, color=c, bold=True)


def slide_subtitle(slide, text, *, y=1.95, color=ITALIC_GRAY, size=20):
    _text(slide, 0.85, y, 12, 0.55, text,
          font=FONT_BODY, size=size, color=color, italic=True)


def page_number(slide, n, total):
    _text(slide, 11.6, 6.95, 1.5, 0.35, f"{n} / {total}",
          font=FONT_BODY, size=14, color=MUTED, align=PP_ALIGN.RIGHT)


def pastel_card(slide, x, y, w, h, *, accent, fill,
                title=None, title_size=22):
    _rounded(slide, x, y, w, h, fill, line_color=accent, line_w=2)
    if title:
        _text(slide, x + 0.25, y + 0.18, w - 0.5, 0.55, title,
              font=FONT_TITLE, size=title_size, color=accent, bold=True)


def circle_number(slide, cx, cy, num_char, accent, r=0.30):
    """Outlined circle with a numeric/figure character inside (e.g. ①②③)."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        _I(cx - r), _I(cy - r), _I(r * 2), _I(r * 2)
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = BG_WHITE
    sh.line.color.rgb = accent
    sh.line.width = _P(2.5)
    sh.shadow.inherit = False
    _text(slide, cx - r, cy - r, r * 2, r * 2, num_char,
          font=FONT_TITLE, size=26, color=accent, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def code_block(slide, x, y, w, h, lines, *, size=12):
    """Dark code block; each line is str or (text, color)."""
    _rect(slide, x, y, w, h, CODE_BG)
    paragraphs = []
    for line in lines:
        if isinstance(line, tuple):
            text, color = line
        else:
            text, color = line, CODE_FG
        paragraphs.append({
            "text": text if text else " ",
            "font": FONT_CODE,
            "size": size,
            "color": color,
            "space_after": 1,
        })
    _multi(slide, x + 0.32, y + 0.22, w - 0.55, h - 0.42, paragraphs)


def callout_box(slide, x, y, w, h, text, *, accent=BLUE, fill=BLUE_PASTEL,
                icon="▶", size=16):
    _rounded(slide, x, y, w, h, fill, line_color=accent, line_w=2)
    _multi(slide, x + 0.3, y, w - 0.6, h, [{
        "text": f"{icon}  {text}",
        "font": FONT_BODY, "size": size, "color": accent, "bold": True,
    }], anchor=MSO_ANCHOR.MIDDLE)


def make_presentation():
    """Create a fresh presentation with our standard 10 × 5.625 dims."""
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


# Map of accent → pastel fill for convenience
PASTEL = {
    bytes(VIOLET):       VIOLET_PASTEL,
    bytes(ORANGE):       ORANGE_PASTEL,
    bytes(TEAL):         TEAL_PASTEL,
    bytes(TEAL_DEEP):    TEAL_PASTEL,
    bytes(PINK):         PINK_PASTEL,
    bytes(BLUE):         BLUE_PASTEL,
    bytes(VIOLET_DEEP):  VIOLET_PASTEL,
}


def pastel_for(accent):
    return PASTEL.get(bytes(accent), VIOLET_PASTEL)
