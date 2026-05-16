#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Build 02-how-mcp-works.pptx from scratch in the new visual style
(matches tool_calling_intro_slides.pptx — violet primary + pastel cards
+ heavy headings + minimal top metadata bar).

Run:
    uv run --with python-pptx python3 tools/build-02-slides.py

Output:
    ./02-how-mcp-works.pptx   (overwritten)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

REPO = Path(__file__).resolve().parent.parent
PPTX = REPO / "02-how-mcp-works.pptx"

# ── Palette ─────────────────────────────────────────────────────────
VIOLET       = RGBColor(0x7B, 0x5C, 0xF5)
VIOLET_DEEP  = RGBColor(0x5B, 0x3E, 0xD9)
ORANGE       = RGBColor(0xFB, 0x70, 0x48)
TEAL         = RGBColor(0x0E, 0xA5, 0x94)
TEAL_DEEP    = RGBColor(0x0F, 0x76, 0x6E)
PINK         = RGBColor(0xEF, 0x44, 0x44)
PINK_DEEP    = RGBColor(0xB9, 0x1C, 0x1C)
BLUE         = RGBColor(0x3B, 0x82, 0xF6)

# Pastel card fills
VIOLET_PASTEL = RGBColor(0xEE, 0xE8, 0xFF)
ORANGE_PASTEL = RGBColor(0xFF, 0xE7, 0xDB)
TEAL_PASTEL   = RGBColor(0xD4, 0xF7, 0xF0)
PINK_PASTEL   = RGBColor(0xFE, 0xE2, 0xE2)
BLUE_PASTEL   = RGBColor(0xDB, 0xEA, 0xFE)
SLATE_PASTEL  = RGBColor(0xF1, 0xF5, 0xF9)

# Neutrals
BG_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MIDNIGHT   = RGBColor(0x0F, 0x14, 0x29)        # for finale slide
INK        = RGBColor(0x1F, 0x29, 0x37)         # body text
INK_SOFT   = RGBColor(0x37, 0x41, 0x51)
MUTED      = RGBColor(0x9C, 0xA3, 0xAF)         # secondary
ITALIC_GRAY = RGBColor(0x6B, 0x72, 0x80)        # italic subtitles
HAIRLINE   = RGBColor(0xE5, 0xE7, 0xEB)         # subtle borders

# Code
CODE_BG       = RGBColor(0x1A, 0x21, 0x38)
CODE_FG       = RGBColor(0xE5, 0xE7, 0xEB)
CODE_STRING   = RGBColor(0xA8, 0xE1, 0xB4)
CODE_KEYWORD  = RGBColor(0xC8, 0xB6, 0xFF)
CODE_COMMENT  = RGBColor(0x9C, 0xA3, 0xAF)
CODE_ORANGE   = RGBColor(0xFB, 0x92, 0x4D)

# ── Fonts ──────────────────────────────────────────────────────────
FONT_TITLE = "Arial Black"
FONT_BODY  = "Calibri"
FONT_CODE  = "Consolas"

# Slide dim — pptx is 10 × 5.625 inches; logical 13.333 × 7.5 scaled by S
SLIDE_W, SLIDE_H = 10.0, 5.625
S = 0.75


def _I(v): return Inches(v * S)
def _P(v): return Pt(v * S)


# ── Low-level primitives ───────────────────────────────────────────
def _rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _I(x), _I(y), _I(w), _I(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = _P(line_w)
    sh.shadow.inherit = False
    return sh


def _rounded(slide, x, y, w, h, color, line_color=None, line_w=None, adj=0.10):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                _I(x), _I(y), _I(w), _I(h))
    sh.adjustments[0] = adj
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = _P(line_w)
    sh.shadow.inherit = False
    return sh


def _oval(slide, x, y, r, color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                _I(x - r), _I(y - r), _I(r * 2), _I(r * 2))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
    sh.shadow.inherit = False
    return sh


def _text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18,
          color=INK, bold=False, italic=False,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = _I(0.04)
    tf.margin_top = tf.margin_bottom = _I(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = _P(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _multi(slide, x, y, w, h, paragraphs, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = _I(0.04)
    tf.margin_top = tf.margin_bottom = _I(0.02)
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if "space_after" in spec:
            p.space_after = _P(spec["space_after"])
        r = p.add_run()
        r.text = spec["text"]
        r.font.name = spec.get("font", FONT_BODY)
        r.font.size = _P(spec.get("size", 18))
        r.font.color.rgb = spec.get("color", INK)
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
    return tb


def _blank_slide(prs, bg=BG_WHITE):
    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    s = prs.slides.add_slide(layout)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                              prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    sp_tree = s.shapes._spTree
    sp_tree.remove(rect._element); sp_tree.insert(2, rect._element)
    return s


# ── Mid-level patterns ─────────────────────────────────────────────
def metadata_bar(slide, num, label, accent=VIOLET):
    """Top-left small colored dot + 'NUM · LABEL' uppercase letter-spaced."""
    _oval(slide, 0.85, 0.55, 0.20, accent)
    _text(slide, 1.2, 0.30, 12, 0.55,
          f"{num}   {label}",
          font=FONT_BODY, size=15, color=MUTED, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)


def slide_title(slide, text, *, y=0.95, size=42, color=INK, accent=None):
    """Large heavy title. If accent given, paints title in accent color."""
    c = accent if accent else color
    _text(slide, 0.85, y, 12, 0.95, text,
          font=FONT_TITLE, size=size, color=c, bold=True)


def slide_subtitle(slide, text, *, y=1.95, color=ITALIC_GRAY, size=20):
    """Italic subtitle below title."""
    _text(slide, 0.85, y, 12, 0.55, text,
          font=FONT_BODY, size=size, color=color, italic=True)


def page_number(slide, n, total):
    _text(slide, 11.6, 6.95, 1.5, 0.35, f"{n} / {total}",
          font=FONT_BODY, size=14, color=MUTED, align=PP_ALIGN.RIGHT)


def pastel_card(slide, x, y, w, h, *, accent, fill,
                title=None, title_size=22):
    """Rounded pastel card with accent border + accent-colored title."""
    _rounded(slide, x, y, w, h, fill, line_color=accent, line_w=2)
    if title:
        _text(slide, x + 0.25, y + 0.18, w - 0.5, 0.55, title,
              font=FONT_TITLE, size=title_size, color=accent, bold=True)


def circle_number(slide, cx, cy, num_char, accent, r=0.30):
    """Outlined circle with a numeric/figure character inside (e.g. ①②③)."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        _I(cx - r), _I(cy - r), _I(r * 2), _I(r * 2)
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = BG_WHITE
    sh.line.color.rgb = accent
    sh.line.width = _P(2.5)
    sh.shadow.inherit = False
    _text(slide, cx - r, cy - r, r * 2, r * 2, num_char,
          font=FONT_TITLE, size=26, color=accent, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def code_block(slide, x, y, w, h, lines, *, size=12):
    """Dark code block; each line is str or (text, color)."""
    _rect(slide, x, y, w, h, CODE_BG)
    paragraphs = []
    for line in lines:
        if isinstance(line, tuple):
            text, color = line
        else:
            text, color = line, CODE_FG
        paragraphs.append({
            "text": text if text else " ",
            "font": FONT_CODE,
            "size": size,
            "color": color,
            "space_after": 1,
        })
    _multi(slide, x + 0.32, y + 0.22, w - 0.55, h - 0.42, paragraphs)


def callout_box(slide, x, y, w, h, text, *, accent=BLUE, fill=BLUE_PASTEL,
                icon="✦", size=16):
    """Bottom-anchored callout with icon + text in pastel rounded box."""
    _rounded(slide, x, y, w, h, fill, line_color=accent, line_w=2)
    _multi(slide, x + 0.3, y, w - 0.6, h, [{
        "text": f"{icon}  {text}",
        "font": FONT_BODY, "size": size, "color": accent, "bold": True,
    }], anchor=MSO_ANCHOR.MIDDLE)


# ── Slide builders ─────────────────────────────────────────────────
TOTAL = 24


def build_cover(prs):
    s = _blank_slide(prs, BG_WHITE)
    # Tiny violet bar top-left
    _rect(s, 0.85, 0.55, 0.55, 0.07, VIOLET)
    _text(s, 0.85, 0.75, 12, 0.4,
          "MCP 入門工作坊  ·  第二講",
          font=FONT_BODY, size=15, color=MUTED, bold=True)

    # Big title
    _text(s, 0.85, 2.0, 12, 1.6,
          "How MCP Works",
          font=FONT_TITLE, size=88, color=INK, bold=True)
    _text(s, 0.85, 3.5, 12, 0.7,
          "把 LLM 接到真實工具的標準協定",
          font=FONT_BODY, size=28, color=VIOLET, bold=True)

    # Author + meta
    _text(s, 0.85, 6.3, 8, 0.4,
          "范耀中  Yao-Chung Fan",
          font=FONT_BODY, size=18, color=INK, bold=True)
    _text(s, 0.85, 6.75, 8, 0.4,
          "國立中興大學  ·  AI 學伴系統實務案例",
          font=FONT_BODY, size=14, color=MUTED)

    # Right-bottom callout
    pastel_card(s, 8.8, 6.0, 4.3, 1.1,
                accent=VIOLET, fill=VIOLET_PASTEL)
    _text(s, 8.95, 6.05, 4.0, 0.5,
          "LLM 不執行任何東西",
          font=FONT_TITLE, size=16, color=VIOLET_DEEP, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, 8.95, 6.5, 4.0, 0.5,
          "所有「執行」都發生在 harness 程式裡",
          font=FONT_BODY, size=13, color=INK, align=PP_ALIGN.CENTER)


def build_agenda(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "00", "A G E N D A", accent=VIOLET)
    slide_title(s, "本講四個段落 + 附錄", y=0.95, size=40)
    slide_subtitle(s, "從 function calling 底層 → MCP 標準化 → 整體架構", y=1.95)

    items = [
        ("①", "Function Calling 怎麼運作", "LLM 只吐字串\nharness 才是執行者",
         ORANGE, ORANGE_PASTEL),
        ("②", "從一個查詢開始",            "場景 → spawn → 兩階段握手\n→ 工具呼叫",
         VIOLET, VIOLET_PASTEL),
        ("③", "Tool 與 Client 機制",       "JSON Schema · description\n注入 LLM 的 tools 參數",
         TEAL, TEAL_PASTEL),
        ("④", "整體架構與附錄",            "四層分工 · 一句話總結\n+ REST vs JSON-RPC",
         PINK, PINK_PASTEL),
    ]
    card_w, card_h = 2.95, 3.5
    gap = 0.20
    x0 = 0.85
    for i, (num, title, body, accent, fill) in enumerate(items):
        x = x0 + i * (card_w + gap)
        _rounded(s, x, 2.85, card_w, card_h, fill,
                 line_color=accent, line_w=2)
        circle_number(s, x + 0.5, 3.25, num, accent, r=0.32)
        _text(s, x + 0.25, 3.95, card_w - 0.5, 0.85, title,
              font=FONT_TITLE, size=18, color=INK, bold=True)
        _multi(s, x + 0.25, 4.85, card_w - 0.5, card_h - 2.0,
               [{"text": line, "font": FONT_BODY, "size": 14,
                 "color": INK_SOFT, "space_after": 4}
                for line in body.split("\n")])
    page_number(s, 2, TOTAL)


# ── Section 01: Function Calling 怎麼運作 ──────────────────────────
def build_fc_intro(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "01", "F U N C T I O N   C A L L I N G", accent=ORANGE)
    slide_title(s, "用一個具體例子搞懂", y=0.95)
    slide_subtitle(s, "—— 從使用者輸入 一路追到 真正執行 的程式碼", y=1.95)

    # User query bubble (large center)
    bub_x, bub_y, bub_w, bub_h = 2.5, 2.9, 8.0, 1.1
    _text(s, bub_x, bub_y - 0.45, bub_w, 0.35, "使用者輸入",
          font=FONT_BODY, size=14, color=MUTED, align=PP_ALIGN.CENTER)
    _rounded(s, bub_x, bub_y, bub_w, bub_h, VIOLET,
             line_color=VIOLET, line_w=0)
    _text(s, bub_x, bub_y, bub_w, bub_h,
          "我家目錄底下有幾個 .py 檔？",
          font=FONT_TITLE, size=26, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Big punchline
    _text(s, 0.85, 4.6, 12, 0.8,
          "LLM 從頭到尾只做一件事 —— 吐字串",
          font=FONT_TITLE, size=32, color=ORANGE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    callout_box(s, 1.5, 5.7, 10.3, 0.75,
                "所有「真的執行」都發生在外面的程式 —— 我們叫它 harness",
                accent=VIOLET, fill=VIOLET_PASTEL, icon="▶")
    page_number(s, 3, TOTAL)


def build_fc_stage12(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "01 · ①", "F U N C T I O N   C A L L I N G   ·   S T A G E   1 + 2",
                 accent=ORANGE)
    slide_title(s, "Stage 1+2:開發者準備 schema,送出 API call", y=0.95, size=30)
    slide_subtitle(s, "tools 只是個 Python dict —— 最終會被 SDK 序列化成 JSON", y=1.85)

    code_lines = [
        ("# 開發者寫的 Python",                                            CODE_COMMENT),
        ("import anthropic",                                              CODE_FG),
        ("client = anthropic.Anthropic()",                                CODE_FG),
        ("",                                                              CODE_FG),
        ("tools = [{",                                                    CODE_FG),
        ('  "name": "execute_bash",',                                     CODE_STRING),
        ('  "description": "Execute a bash command. Returns stdout.",',  CODE_STRING),
        ('  "input_schema": {',                                           CODE_FG),
        ('    "type": "object",',                                         CODE_FG),
        ('    "properties": {"command": {"type": "string"}},',           CODE_FG),
        ('    "required": ["command"]',                                   CODE_FG),
        ("  }",                                                           CODE_FG),
        ("}]",                                                            CODE_FG),
        ("",                                                              CODE_FG),
        ("response = client.messages.create(",                            CODE_FG),
        ('    model="claude-opus-4-7",',                                  CODE_FG),
        ("    tools=tools,        ← 把工具 schema 一起送",                 CODE_ORANGE),
        ("    messages=[{'role':'user',",                                 CODE_FG),
        ("               'content':'我家目錄底下有幾個 .py 檔？'}]",       CODE_FG),
        (")",                                                             CODE_FG),
    ]
    code_block(s, 0.85, 2.5, 7.4, 4.0, code_lines, size=10)

    # Right side callouts
    pastel_card(s, 8.5, 2.5, 4.4, 1.7, accent=VIOLET, fill=VIOLET_PASTEL,
                title="關鍵觀察")
    _multi(s, 8.75, 3.05, 4.0, 1.2, [
        {"text": "tools 只是一個 dict。",
         "font": FONT_BODY, "size": 14, "color": INK,
         "bold": True, "space_after": 4},
        {"text": "它沒有任何「綁定的函式指標」。",
         "font": FONT_BODY, "size": 13, "color": INK, "space_after": 0},
    ])

    pastel_card(s, 8.5, 4.4, 4.4, 2.1, accent=TEAL, fill=TEAL_PASTEL,
                title="模型看到什麼")
    _multi(s, 8.75, 4.95, 4.0, 1.6, [
        {"text": "「有一個叫 execute_bash 的東西,",
         "font": FONT_BODY, "size": 13, "color": INK, "space_after": 3},
        {"text": "  描述是這樣,參數長這樣 …」",
         "font": FONT_BODY, "size": 13, "color": INK, "space_after": 8},
        {"text": "→ 跟讀到任何自然語言沒差別。",
         "font": FONT_BODY, "size": 13, "color": TEAL_DEEP,
         "italic": True, "space_after": 0},
    ])
    page_number(s, 4, TOTAL)


def build_fc_stage3(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "01 · ②", "F U N C T I O N   C A L L I N G   ·   S T A G E   3",
                 accent=ORANGE)
    slide_title(s, "Stage 3:LLM 回了一段 JSON", y=0.95)
    slide_subtitle(s, "它就是吐了一段字串 —— 沒有「呼叫」任何東西", y=1.85)

    json_lines = [
        ("{",                                                  CODE_FG),
        ('  "id": "msg_abc123",',                              CODE_FG),
        ('  "stop_reason": "tool_use",       ← 旗號',           CODE_ORANGE),
        ('  "content": [',                                     CODE_FG),
        ('    { "type": "text",',                              CODE_FG),
        ('      "text": "讓我用 find 幫你數一下。" },',         CODE_STRING),
        ('    { "type": "tool_use",',                          CODE_FG),
        ('      "id": "toolu_01XYZ",',                         CODE_FG),
        ('      "name": "execute_bash",',                      CODE_FG),
        ('      "input": {',                                   CODE_FG),
        ('        "command": "find ~ -name \'*.py\' -type f | wc -l"', CODE_STRING),
        ('      }',                                            CODE_FG),
        ('    }',                                              CODE_FG),
        ('  ]',                                                CODE_FG),
        ('}',                                                  CODE_FG),
    ]
    code_block(s, 0.85, 2.5, 8.0, 3.7, json_lines, size=11)

    pastel_card(s, 9.1, 2.5, 3.8, 1.8, accent=ORANGE, fill=ORANGE_PASTEL,
                title="stop_reason")
    _multi(s, 9.3, 3.05, 3.5, 1.4, [
        {"text": "tool_use ←→ end_turn",
         "font": FONT_CODE, "size": 12, "color": ORANGE,
         "bold": True, "space_after": 6},
        {"text": "前者:harness 接手執行。\n後者:迴圈結束,印出。",
         "font": FONT_BODY, "size": 12, "color": INK, "space_after": 0},
    ])

    pastel_card(s, 9.1, 4.5, 3.8, 1.7, accent=TEAL, fill=TEAL_PASTEL,
                title="command 這串字")
    _multi(s, 9.3, 5.05, 3.5, 1.2, [
        {"text": "是 LLM token-by-token 生出來的。",
         "font": FONT_BODY, "size": 12, "color": INK, "space_after": 4},
        {"text": "Constrained decoding 確保 JSON 一定合法。",
         "font": FONT_BODY, "size": 12, "color": TEAL_DEEP,
         "italic": True, "space_after": 0},
    ])

    # Punchline at bottom
    _rounded(s, 0.85, 6.45, 12, 0.6, ORANGE_PASTEL,
             line_color=ORANGE, line_w=2)
    _text(s, 0.85, 6.45, 12, 0.6,
          "LLM 從來沒有「呼叫」過任何工具 —— 它就是吐了一段 JSON 字串",
          font=FONT_TITLE, size=18, color=ORANGE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    page_number(s, 5, TOTAL)


def build_fc_stage4(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "01 · ③", "F U N C T I O N   C A L L I N G   ·   S T A G E   4",
                 accent=PINK)
    slide_title(s, "Stage 4:Harness 真的執行", y=0.95, accent=PINK)
    slide_subtitle(s, "subprocess.run() —— 字串穿越文字宇宙、進入真實電腦的那一刻", y=1.85)

    code_lines = [
        ("# Harness 拿到 LLM 的 tool_use 後", CODE_COMMENT),
        ("for block in response.content:",   CODE_FG),
        ('    if block.type == "tool_use":', CODE_FG),
        ('        cmd = block.input["command"]', CODE_FG),
        ("",                                CODE_FG),
        ("        # ↓↓↓  這一行才是真的執行 ↓↓↓", CODE_COMMENT),
        ("        result = subprocess.run(", CODE_ORANGE),
        ("            cmd,",                CODE_FG),
        ("            shell=True,",         CODE_ORANGE),
        ("            capture_output=True,", CODE_FG),
        ("            text=True,",          CODE_FG),
        ("            timeout=30,",         CODE_FG),
        ("        )",                       CODE_FG),
        ('        tool_output = result.stdout.strip()  # "42"', CODE_FG),
    ]
    code_block(s, 0.85, 2.5, 7.5, 3.5, code_lines, size=11)

    pastel_card(s, 8.6, 2.5, 4.3, 1.8, accent=PINK, fill=PINK_PASTEL,
                title="信任邊界")
    _multi(s, 8.85, 3.05, 4.0, 1.4, [
        {"text": "subprocess.run(cmd, shell=True)",
         "font": FONT_CODE, "size": 11, "color": PINK, "bold": True,
         "space_after": 4},
        {"text": "把 LLM 的字串交給 OS shell。從這一刻起,跟你親手打 100% 一樣。",
         "font": FONT_BODY, "size": 12, "color": INK, "space_after": 0},
    ])

    pastel_card(s, 8.6, 4.5, 4.3, 1.5, accent=TEAL_DEEP, fill=SLATE_PASTEL,
                title="OS 在這一刻做的事")
    _multi(s, 8.85, 5.05, 4.0, 1.0, [
        {"text": "/bin/bash 解析字串",
         "font": FONT_BODY, "size": 12, "color": INK, "space_after": 3},
        {"text": "→ fork 出 find / wc processes",
         "font": FONT_BODY, "size": 12, "color": INK, "space_after": 3},
        {"text": "→ stdout 回傳 \"42\"",
         "font": FONT_CODE, "size": 12, "color": TEAL_DEEP,
         "bold": True, "space_after": 0},
    ])

    # Bottom warning
    _rounded(s, 0.85, 6.25, 12, 0.85, PINK_PASTEL,
             line_color=PINK, line_w=2)
    _multi(s, 1.1, 6.3, 11.5, 0.8, [
        {"text": "△   subprocess.run(LLM 輸出) 本質上就是 eval(LLM 輸出)",
         "font": FONT_BODY, "size": 15, "color": PINK_DEEP, "bold": True,
         "space_after": 2},
        {"text": "若 prompt injection 攻擊,LLM 會生出 rm -rf ~,harness 不長眼就會照跑。",
         "font": FONT_BODY, "size": 12, "color": INK, "space_after": 0},
    ])
    page_number(s, 6, TOTAL)


def build_fc_stage56(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "01 · ④", "F U N C T I O N   C A L L I N G   ·   S T A G E   5 + 6",
                 accent=ORANGE)
    slide_title(s, "Stage 5+6:結果塞回,LLM 給最終答案", y=0.95, size=32)
    slide_subtitle(s, "messages 陣列繼續滾動,直到 stop_reason: end_turn", y=1.85)

    msg_lines = [
        ("[user]",                                                BLUE),
        ('    "我家目錄底下有幾個 .py 檔？"',                       CODE_FG),
        ("",                                                       CODE_FG),
        ("[assistant]",                                            VIOLET),
        ('    "讓我用 find 幫你數一下。"',                          CODE_STRING),
        ("    + tool_use { execute_bash, find ~ ... }",            CODE_FG),
        ("",                                                       CODE_FG),
        ("[user]   ← 注意:role 是 user",                          TEAL),
        ('    tool_result { content="42" }',                       CODE_FG),
        ("",                                                       CODE_FG),
        ("[assistant]   stop_reason: end_turn",                    ORANGE),
        ('    "你家目錄底下總共有 42 個 .py 檔。"',                  CODE_STRING),
    ]
    code_block(s, 0.85, 2.5, 7.5, 4.0, msg_lines, size=12)

    pastel_card(s, 8.6, 2.5, 4.3, 1.9, accent=TEAL, fill=TEAL_PASTEL,
                title="tool_result 用 user 送回")
    _text(s, 8.85, 3.05, 4.0, 1.5,
          "從 model 角度,它中斷自己問了外面世界一個問題,外面回了一句。現在輪它繼續講。",
          font=FONT_BODY, size=12, color=INK)

    pastel_card(s, 8.6, 4.6, 4.3, 1.9, accent=ORANGE, fill=ORANGE_PASTEL,
                title="完整 loop")
    _multi(s, 8.85, 5.15, 4.0, 1.5, [
        {"text": "tool_use",
         "font": FONT_CODE, "size": 13, "color": VIOLET, "bold": True,
         "space_after": 2},
        {"text": "→ run → tool_result",
         "font": FONT_CODE, "size": 13, "color": INK, "space_after": 2},
        {"text": "→ … → end_turn",
         "font": FONT_CODE, "size": 13, "color": ORANGE, "bold": True,
         "space_after": 0},
    ])

    _text(s, 0.85, 6.7, 12, 0.4,
          "可能跑很多輪 —— LLM 可以連續呼叫多個工具,直到它覺得夠了。",
          font=FONT_BODY, size=14, color=MUTED, italic=True,
          align=PP_ALIGN.CENTER)
    page_number(s, 7, TOTAL)


def build_fc_takeaway(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "01 · ⑤", "F U N C T I O N   C A L L I N G   ·   T A K E A W A Y",
                 accent=VIOLET)
    slide_title(s, "LLM ↔ Harness 職責分工", y=0.95)
    slide_subtitle(s, "MCP 在這個機制上加了什麼?把 harness ↔ tool 標準化、抽到獨立 process", y=1.85)

    # Left: LLM card
    pastel_card(s, 0.85, 2.7, 5.9, 3.5, accent=VIOLET, fill=VIOLET_PASTEL,
                title="LLM 做的事")
    _multi(s, 1.15, 3.25, 5.4, 2.9, [
        {"text": "✓  生成 tool_use JSON",
         "font": FONT_BODY, "size": 17, "color": INK, "bold": True,
         "space_after": 8},
        {"text": "✓  讀 tool schema,生 input JSON",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 8},
        {"text": "✓  Constrained decoding 確保 JSON 合法",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 8},
        {"text": "✗  不知道工具真正做什麼",
         "font": FONT_BODY, "size": 16, "color": MUTED, "space_after": 8},
        {"text": "✗  碰不到 OS、檔案、網路",
         "font": FONT_BODY, "size": 16, "color": PINK, "bold": True,
         "space_after": 0},
    ])

    # Right: Harness card
    pastel_card(s, 7.05, 2.7, 5.85, 3.5, accent=TEAL, fill=TEAL_PASTEL,
                title="Harness 做的事")
    _multi(s, 7.35, 3.25, 5.35, 2.9, [
        {"text": "✓  解析 tool_use,真的執行",
         "font": FONT_BODY, "size": 17, "color": INK, "bold": True,
         "space_after": 8},
        {"text": "✓  維護 messages 陣列",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 8},
        {"text": "✓  Loop 直到 stop_reason: end_turn",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 8},
        {"text": "✓  處理 timeout / 錯誤 / 結果格式",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 8},
        {"text": "✓  加 sandbox / 安全防護",
         "font": FONT_BODY, "size": 16, "color": TEAL_DEEP, "bold": True,
         "space_after": 0},
    ])

    callout_box(s, 0.85, 6.4, 12, 0.65,
                "MCP = 把 harness ↔ tool 的協定標準化,抽到獨立 process",
                accent=VIOLET, fill=VIOLET_PASTEL, icon="▶", size=15)
    page_number(s, 8, TOTAL)


# ── Section 02: 從一個查詢開始 ─────────────────────────────────────
def build_scenario(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "02", "F R O M   A   U S E R   Q U E R Y", accent=VIOLET)
    slide_title(s, "場景:使用者問了一個問題", y=0.95)
    slide_subtitle(s, "LLM 判斷該呼叫哪個工具 —— 但它自己不會呼叫", y=1.85)

    # User bubble top-right
    _text(s, 6.5, 2.45, 6, 0.35, "使用者", font=FONT_BODY,
          size=14, color=MUTED, align=PP_ALIGN.RIGHT)
    _rounded(s, 6.5, 2.85, 6.4, 0.95, BLUE,
             line_color=BLUE, line_w=0)
    _text(s, 6.5, 2.85, 6.4, 0.95,
          "幫我查最近圖書館有什麼新書",
          font=FONT_TITLE, size=20, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # LLM thinking box bottom-left
    pastel_card(s, 0.85, 4.4, 7.0, 1.85,
                accent=VIOLET, fill=VIOLET_PASTEL,
                title="LLM (Sonnet) 內心戲")
    _multi(s, 1.1, 4.95, 6.6, 1.3, [
        {"text": "「嗯…需要查新書資料庫」",
         "font": FONT_BODY, "size": 18, "color": INK, "space_after": 4},
        {"text": "→  呼叫 search_new_books 工具",
         "font": FONT_CODE, "size": 18, "color": VIOLET, "bold": True,
         "space_after": 0},
    ])

    callout_box(s, 0.85, 6.5, 12, 0.55,
                "但 LLM 自己不會呼叫 —— 要靠 Parent 程序去跟 Child 程序講話",
                accent=PINK, fill=PINK_PASTEL, icon="!", size=15)
    page_number(s, 9, TOTAL)


def build_act1(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "02 · ①", "F R O M   A   U S E R   Q U E R Y   ·   A C T   1",
                 accent=VIOLET)
    slide_title(s, "Act 1:Parent 開出 Child", y=0.95)
    slide_subtitle(s, "spawn 子程序 + 兩條 stdio pipe = 雙向通訊管道", y=1.85)

    # Parent box (left)
    parent_x, parent_y, parent_w, parent_h = 1.0, 2.9, 4.2, 2.0
    _rounded(s, parent_x, parent_y, parent_w, parent_h, VIOLET_PASTEL,
             line_color=VIOLET, line_w=2.5)
    _text(s, parent_x, parent_y + 0.15, parent_w, 0.55,
          "Parent Process",
          font=FONT_TITLE, size=22, color=VIOLET_DEEP, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, parent_x, parent_y + 0.85, parent_w, 0.4, "Node.js",
          font=FONT_CODE, size=16, color=INK_SOFT,
          align=PP_ALIGN.CENTER)
    # Sonnet badge inside
    _rounded(s, parent_x + 0.85, parent_y + 1.35, parent_w - 1.7, 0.45,
             ORANGE, line_color=ORANGE, line_w=0)
    _text(s, parent_x + 0.85, parent_y + 1.35, parent_w - 1.7, 0.45,
          "LLM · Sonnet",
          font=FONT_CODE, size=13, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Child box (right)
    child_x, child_y, child_w, child_h = 8.1, 2.9, 4.2, 2.0
    _rounded(s, child_x, child_y, child_w, child_h, TEAL_PASTEL,
             line_color=TEAL, line_w=2.5)
    _text(s, child_x, child_y + 0.15, child_w, 0.55,
          "Child Process",
          font=FONT_TITLE, size=22, color=TEAL_DEEP, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, child_x, child_y + 1.0, child_w, 0.4, "Python MCP Server",
          font=FONT_CODE, size=16, color=INK_SOFT,
          align=PP_ALIGN.CENTER)

    # Pipes between
    _rounded(s, 5.4, 3.4, 2.5, 0.3, VIOLET_PASTEL,
             line_color=VIOLET, line_w=1.5)
    _text(s, 5.4, 3.4, 2.5, 0.3, "stdout",
          font=FONT_CODE, size=11, color=VIOLET_DEEP,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _rounded(s, 5.4, 4.1, 2.5, 0.3, TEAL_PASTEL,
             line_color=TEAL, line_w=1.5)
    _text(s, 5.4, 4.1, 2.5, 0.3, "stdin",
          font=FONT_CODE, size=11, color=TEAL_DEEP,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom 3 numbered points
    _multi(s, 0.85, 5.3, 12, 1.8, [
        {"text": "①  Parent (Node.js) 用 spawn() 開出 Child (Python) 子程序",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 5},
        {"text": "②  LLM (Sonnet) 是 Parent 內的 client component —— 不會自己出去呼叫工具",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 5},
        {"text": "③  兩條 stdio pipe → Parent ⇄ Child 可以雙向講話",
         "font": FONT_BODY, "size": 16, "color": INK, "space_after": 0},
    ])
    page_number(s, 10, TOTAL)


def build_act2(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "02 · ②", "F R O M   A   U S E R   Q U E R Y   ·   A C T   2",
                 accent=VIOLET)
    slide_title(s, "Act 2:兩階段握手", y=0.95)
    slide_subtitle(s, "互換能力 → 拿到工具清單,兩步都過才算真的連上", y=1.85)

    # Left: process flow
    pastel_card(s, 0.85, 2.7, 6.4, 4.0, accent=VIOLET, fill=BG_WHITE,
                title="訊息往返")
    _multi(s, 1.15, 3.3, 6.0, 3.3, [
        {"text": "①   Parent  →  Child",
         "font": FONT_TITLE, "size": 16, "color": VIOLET, "bold": True,
         "space_after": 2},
        {"text": "      「我來了」  +  「給我能力清單」",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 12},

        {"text": "②   Child  →  Parent",
         "font": FONT_TITLE, "size": 16, "color": TEAL, "bold": True,
         "space_after": 2},
        {"text": "      「我能做什麼」  +  「8 個工具清單」",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 12},

        {"text": "→  雙向 JSON-RPC 訊息往返",
         "font": FONT_BODY, "size": 14, "color": ITALIC_GRAY, "italic": True,
         "space_after": 0},
    ])

    # Right: checklist
    pastel_card(s, 7.5, 2.7, 5.4, 4.0, accent=TEAL, fill=TEAL_PASTEL,
                title="握手進度")
    # Step 1 — checked
    _rounded(s, 7.85, 3.5, 0.45, 0.45, TEAL, line_color=TEAL, line_w=0)
    _text(s, 7.85, 3.5, 0.45, 0.45, "✓",
          font=FONT_TITLE, size=20, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 8.5, 3.5, 4.4, 0.45,
          "Step 1:打過招呼",
          font=FONT_TITLE, size=17, color=TEAL_DEEP, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)

    # Step 2 — checked
    _rounded(s, 7.85, 4.4, 0.45, 0.45, TEAL, line_color=TEAL, line_w=0)
    _text(s, 7.85, 4.4, 0.45, 0.45, "✓",
          font=FONT_TITLE, size=20, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 8.5, 4.4, 4.4, 0.45,
          "Step 2:拿到工具清單",
          font=FONT_TITLE, size=17, color=TEAL_DEEP, bold=True,
          anchor=MSO_ANCHOR.MIDDLE)

    _text(s, 7.85, 5.7, 5.0, 0.7,
          "兩個都打勾 才算真的連上",
          font=FONT_BODY, size=15, color=INK, italic=True,
          anchor=MSO_ANCHOR.MIDDLE)

    callout_box(s, 0.85, 6.85, 12, 0.0,  # spacer
                "", accent=VIOLET, fill=BG_WHITE, icon="")
    _text(s, 0.85, 6.85, 12, 0.4,
          "中間時刻 (Step 1 ✓ 但 Step 2 ✗) Parent 還不知道工具清單,送 tool call 也沒用。",
          font=FONT_BODY, size=13, color=MUTED, italic=True,
          align=PP_ALIGN.CENTER)
    page_number(s, 11, TOTAL)


def build_act3(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "02 · ③", "F R O M   A   U S E R   Q U E R Y   ·   A C T   3",
                 accent=VIOLET)
    slide_title(s, "Act 3:工具呼叫", y=0.95)
    slide_subtitle(s, "LLM 決定 → token 飛去 Child → 結果飛回 → 找對應請求", y=1.85)

    # Flow as 5 numbered steps in a horizontal strip
    items = [
        ("①", "LLM 決定", "search_new_books(\"新書\")",                 VIOLET),
        ("②", "Parent 記下", "等候第 4 個請求 · 30s 內",                 ORANGE),
        ("③", "Child 執行", "Python MCP Server 跑 search_new_books()", TEAL),
        ("④", "結果飛回", "10 筆新書",                                   TEAL_DEEP),
        ("⑤", "Parent 回呼", "找到對應請求 → 回給使用者",                 VIOLET),
    ]
    card_w = 2.40
    gap = 0.10
    x0 = 0.55
    for i, (num, title, body, accent) in enumerate(items):
        x = x0 + i * (card_w + gap)
        fill = {VIOLET: VIOLET_PASTEL, ORANGE: ORANGE_PASTEL,
                TEAL: TEAL_PASTEL, TEAL_DEEP: TEAL_PASTEL}.get(accent, VIOLET_PASTEL)
        _rounded(s, x, 2.8, card_w, 3.0, fill,
                 line_color=accent, line_w=2)
        circle_number(s, x + 0.45, 3.20, num, accent, r=0.28)
        _text(s, x + 0.15, 3.7, card_w - 0.3, 0.5, title,
              font=FONT_TITLE, size=15, color=INK, bold=True,
              align=PP_ALIGN.CENTER)
        _text(s, x + 0.15, 4.3, card_w - 0.3, 1.4, body,
              font=FONT_BODY, size=12, color=INK_SOFT,
              align=PP_ALIGN.CENTER)

    callout_box(s, 0.85, 6.25, 12, 0.65,
                "「第 4 個請求」是 JSON-RPC 給每個 request 的編號,讓 response 對得回來",
                accent=VIOLET, fill=VIOLET_PASTEL, icon="▶", size=14)
    page_number(s, 12, TOTAL)


def build_video_cue(prs):
    s = _blank_slide(prs, MIDNIGHT)
    _text(s, 0.85, 1.8, 12, 0.8, "▶  切換到影片",
          font=FONT_TITLE, size=44, color=ORANGE, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, 0.85, 2.8, 12, 0.5,
          "02-mcp-connection-video.mp4   ·   2 : 1 7",
          font=FONT_CODE, size=22, color=MUTED, align=PP_ALIGN.CENTER)

    pastel_card(s, 2.5, 4.0, 8.3, 2.3,
                accent=VIOLET, fill=RGBColor(0x1E, 0x20, 0x3A))
    _multi(s, 2.85, 4.15, 7.6, 2.1, [
        {"text": "影片重點",
         "font": FONT_TITLE, "size": 16, "color": VIOLET, "bold": True,
         "space_after": 8},
        {"text": "1. 場景開場 — 使用者問「新書」 → LLM 判斷要呼叫工具",
         "font": FONT_BODY, "size": 13, "color": CODE_FG, "space_after": 4},
        {"text": "2. Act 1 — Parent / Child 的 spawn 與 pipe",
         "font": FONT_BODY, "size": 13, "color": CODE_FG, "space_after": 4},
        {"text": "3. Act 2 — 兩階段握手 checkbox 視覺",
         "font": FONT_BODY, "size": 13, "color": CODE_FG, "space_after": 4},
        {"text": "4. Act 3 — token 飛去、等候第 4 個請求、結果飛回",
         "font": FONT_BODY, "size": 13, "color": CODE_FG, "space_after": 4},
        {"text": "5. Frame closure — 「新書」一路回到使用者",
         "font": FONT_BODY, "size": 13, "color": CODE_FG, "space_after": 0},
    ])
    page_number(s, 13, TOTAL)


# ── Section 03: Tool 註冊與描述 ─────────────────────────────────────
def build_tools_list(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "03", "T O O L   ·   註 冊 與 描 述", accent=TEAL)
    slide_title(s, "tools/list:Server 回傳工具定義", y=0.95, size=30)
    slide_subtitle(s, "Client 一啟動就 list,把每個工具的 schema 收進來", y=1.85)

    code_lines = [
        ("# Client → Server",                              CODE_COMMENT),
        ('{ "method": "tools/list", "id": 2 }',            CODE_FG),
        ("",                                               CODE_FG),
        ("# Server → Client  (興大系統共 33 個)",          CODE_COMMENT),
        ('{ "result": { "tools": [{',                      CODE_FG),
        ('    "name": "search_library_books",',            CODE_STRING),
        ('    "description": "搜尋中興大學圖書館的館藏書籍,', CODE_STRING),
        ('                    支援關鍵字查詢",',             CODE_STRING),
        ('    "inputSchema": {',                           CODE_FG),
        ('      "type": "object",',                        CODE_FG),
        ('      "properties": {',                          CODE_FG),
        ('        "keyword": {"type": "string"}',          CODE_FG),
        ('      }, "required": ["keyword"]',               CODE_FG),
        ('    }',                                          CODE_FG),
        ('  }, { "name": "get_department_courses", ... },', CODE_FG),
        ('  …共 33 個',                                     CODE_COMMENT),
        ('] }, "id": 2 }',                                 CODE_FG),
    ]
    code_block(s, 0.85, 2.5, 12, 4.3, code_lines, size=11)
    page_number(s, 14, TOTAL)


def build_three_fields(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "03 · ①", "T O O L   ·   三 個 關 鍵 欄 位", accent=TEAL)
    slide_title(s, "每個 Tool 的三個關鍵欄位", y=0.95)
    slide_subtitle(s, "LLM 看到的就只有這三個,所以 description 決定一切", y=1.85)

    items = [
        ("name", "工具的唯一識別名稱",
         '"search_library_books"',
         "LLM 在 tool_use 回傳時引用這個名稱", VIOLET, VIOLET_PASTEL),
        ("description", "自然語言描述工具的用途",
         '"搜尋中興大學圖書館的館藏書籍"',
         "這是 LLM 判斷「要不要用」的最重要依據", ORANGE, ORANGE_PASTEL),
        ("inputSchema", "JSON Schema 定義參數格式",
         '{ "type": "object",\n  "properties": {…} }',
         "LLM 據此產生正確的 arguments", TEAL, TEAL_PASTEL),
    ]
    card_h = 4.1
    card_w = 3.95
    gap = 0.20
    x0 = 0.55
    for i, (key, desc, example, footnote, accent, fill) in enumerate(items):
        x = x0 + i * (card_w + gap)
        _rounded(s, x, 2.75, card_w, card_h, fill,
                 line_color=accent, line_w=2)
        _text(s, x + 0.2, 2.95, card_w - 0.4, 0.5, key,
              font=FONT_CODE, size=20, color=accent, bold=True)
        _text(s, x + 0.2, 3.55, card_w - 0.4, 0.7, desc,
              font=FONT_BODY, size=13, color=INK)
        # Example block
        _rounded(s, x + 0.2, 4.45, card_w - 0.4, 1.1, CODE_BG,
                 line_color=CODE_BG)
        _text(s, x + 0.3, 4.5, card_w - 0.6, 1.0, example,
              font=FONT_CODE, size=10, color=CODE_STRING,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x + 0.2, 5.75, card_w - 0.4, 1.05, footnote,
              font=FONT_BODY, size=12, color=INK_SOFT, italic=True)
    page_number(s, 15, TOTAL)


def build_desc_quality(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "03 · ②", "T O O L   ·   D E S C R I P T I O N", accent=TEAL)
    slide_title(s, "description 寫得好不好,決定 LLM 選不選得對", y=0.95, size=30)
    slide_subtitle(s, "LLM 從來不會直接接觸 MCP Server,只看 name + description + inputSchema", y=1.85)

    # Bad
    _rounded(s, 0.85, 2.75, 5.95, 3.0, PINK_PASTEL,
             line_color=PINK, line_w=2)
    _text(s, 1.1, 2.9, 5.5, 0.5,
          "✗   模糊的 description",
          font=FONT_TITLE, size=18, color=PINK_DEEP, bold=True)
    _multi(s, 1.1, 3.45, 5.5, 2.2, [
        {"text": '"搜尋課程"',
         "font": FONT_CODE, "size": 17, "color": INK, "space_after": 12},
        {"text": "問題:8 個課程相關工具",
         "font": FONT_BODY, "size": 14, "color": INK, "space_after": 4},
        {"text": "  (關鍵字 / 系所 / 教師 / 類型…)",
         "font": FONT_BODY, "size": 13, "color": INK_SOFT, "space_after": 4},
        {"text": "LLM 分不清楚要用哪一個。",
         "font": FONT_BODY, "size": 13, "color": PINK_DEEP, "italic": True,
         "space_after": 0},
    ])

    # Good
    _rounded(s, 7.0, 2.75, 5.95, 3.0, TEAL_PASTEL,
             line_color=TEAL, line_w=2)
    _text(s, 7.25, 2.9, 5.5, 0.5,
          "✓   精確的 description",
          font=FONT_TITLE, size=18, color=TEAL_DEEP, bold=True)
    _multi(s, 7.25, 3.45, 5.5, 2.2, [
        {"text": '"依系所代碼查詢該系所的',
         "font": FONT_CODE, "size": 13, "color": INK, "space_after": 2},
        {"text": ' 所有課程,回傳課程名稱、',
         "font": FONT_CODE, "size": 13, "color": INK, "space_after": 2},
        {"text": ' 學分數、授課教師"',
         "font": FONT_CODE, "size": 13, "color": INK, "space_after": 10},
        {"text": "LLM 一看就知道:「電機系有哪些課」→ 用這個",
         "font": FONT_BODY, "size": 13, "color": TEAL_DEEP, "italic": True,
         "space_after": 0},
    ])

    # Bottom key insight
    _rounded(s, 0.85, 6.0, 12, 0.95, CODE_BG)
    _multi(s, 1.1, 6.05, 11.5, 0.9, [
        {"text": "  關鍵洞察",
         "font": FONT_TITLE, "size": 13, "color": ORANGE, "bold": True,
         "space_after": 2},
        {"text": "  description 的品質  =  工具被正確使用的機率",
         "font": FONT_BODY, "size": 15, "color": CODE_FG, "bold": True,
         "space_after": 0},
    ])
    page_number(s, 16, TOTAL)


# ── Section 04: Client 整合機制 + 整體架構 ──────────────────────────
def build_client_integration(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "04", "C L I E N T   ·   整 合 機 制", accent=VIOLET)
    slide_title(s, "工具清單 → LLM 的 tools 參數", y=0.95)
    slide_subtitle(s, "Client 把所有 MCP Server 的工具收齊,轉成 LLM API 的 tools 欄位", y=1.85)

    # Flow strip
    arrow_items = [
        ("33 個 MCP\nServers",   VIOLET),
        ("Client 收集",          ORANGE),
        ("轉換",                  TEAL),
        ("LLM API\n呼叫",         PINK),
    ]
    card_w = 2.5
    gap = 0.25
    x0 = 1.05
    for i, (txt, color) in enumerate(arrow_items):
        x = x0 + i * (card_w + gap)
        fill = {VIOLET: VIOLET_PASTEL, ORANGE: ORANGE_PASTEL,
                TEAL: TEAL_PASTEL, PINK: PINK_PASTEL}[color]
        _rounded(s, x, 2.8, card_w, 1.0, fill,
                 line_color=color, line_w=2)
        _text(s, x, 2.8, card_w, 1.0, txt,
              font=FONT_TITLE, size=13, color=color, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(arrow_items) - 1:
            _text(s, x + card_w + 0.02, 2.8, gap - 0.05, 1.0, "→",
                  font=FONT_TITLE, size=22, color=MUTED, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Code block bottom
    code_lines = [
        ("// raw-mcp-client.js  —  注入工具到 Claude API",  CODE_COMMENT),
        ("const response = await anthropic.messages.create({", CODE_FG),
        ('  model: "claude-sonnet-4-20250514",',           CODE_FG),
        ("  messages: conversationHistory,",               CODE_FG),
        ("  tools: allMcpTools.map(t => ({",               CODE_ORANGE),
        ("    name: t.name,             // 來自 MCP Server", CODE_FG),
        ("    description: t.description, // 來自 MCP Server", CODE_FG),
        ("    input_schema: t.inputSchema // 來自 MCP Server", CODE_FG),
        ("  }))",                                          CODE_FG),
        ("});",                                            CODE_FG),
    ]
    code_block(s, 0.85, 4.2, 12, 2.7, code_lines, size=11)
    page_number(s, 17, TOTAL)


def build_llm_view(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "04 · ①", "C L I E N T   ·   L L M 眼中的世界", accent=VIOLET)
    slide_title(s, "LLM 眼中的世界 —— 只看到工具描述", y=0.95, size=30)
    slide_subtitle(s, "Server 怎麼跑、語言、資料庫,LLM 都看不到。它只有文字。", y=1.85)

    # Left card: 看到的
    _rounded(s, 0.85, 2.75, 5.95, 4.1, CODE_BG)
    _text(s, 1.1, 2.9, 5.5, 0.45,
          "LLM 看到的",
          font=FONT_TITLE, size=16, color=ORANGE, bold=True)
    _multi(s, 1.1, 3.45, 5.7, 3.3, [
        {"text": "Available tools:",
         "font": FONT_CODE, "size": 13, "color": CODE_KEYWORD,
         "space_after": 8},
        {"text": "1.  search_library_books",
         "font": FONT_CODE, "size": 13, "color": CODE_FG, "space_after": 2},
        {"text": "      搜尋圖書館館藏書籍",
         "font": FONT_BODY, "size": 11, "color": MUTED, "space_after": 8},
        {"text": "2.  get_department_courses",
         "font": FONT_CODE, "size": 13, "color": CODE_FG, "space_after": 2},
        {"text": "      依系所代碼查詢課程",
         "font": FONT_BODY, "size": 11, "color": MUTED, "space_after": 8},
        {"text": "3.  search_teachers",
         "font": FONT_CODE, "size": 13, "color": CODE_FG, "space_after": 2},
        {"text": "      搜尋教師資訊與研究",
         "font": FONT_BODY, "size": 11, "color": MUTED, "space_after": 10},
        {"text": "…共 33 個",
         "font": FONT_BODY, "size": 12, "color": MUTED, "italic": True,
         "space_after": 0},
    ])

    # Right card: 看不到的
    _rounded(s, 7.0, 2.75, 5.95, 4.1, SLATE_PASTEL,
             line_color=MUTED, line_w=1)
    _text(s, 7.25, 2.9, 5.5, 0.45,
          "LLM 看不到的",
          font=FONT_TITLE, size=16, color=PINK_DEEP, bold=True)
    _multi(s, 7.25, 3.45, 5.5, 3.3, [
        {"text": "✗  Server 怎麼啟動的",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 8},
        {"text": "✗  Server 用什麼語言",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 8},
        {"text": "✗  資料庫連線方式",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 8},
        {"text": "✗  網路拓撲",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 8},
        {"text": "✗  認證機制",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 8},
        {"text": "✗  JSON-RPC 細節",
         "font": FONT_BODY, "size": 15, "color": INK, "space_after": 0},
    ])
    page_number(s, 18, TOTAL)


def build_architecture(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "04 · ②", "M C P   ·   整 體 架 構", accent=VIOLET)
    slide_title(s, "MCP 整體架構:四層分工", y=0.95)
    slide_subtitle(s, "Host → Client → Server → Tool —— 每層只看下一層,LLM 只看文字", y=1.85)

    layers = [
        ("Host",   "應用程式 / 入口",        "興大 AI 學伴 web UI",                          PINK,   PINK_PASTEL),
        ("Client", "包含 LLM + MCP client", "Node.js + Sonnet  (= 影片裡的 Parent)",          VIOLET, VIOLET_PASTEL),
        ("Server", "個別 MCP server 子程序", "Python MCP Server  (= 影片裡的 Child)",          TEAL,   TEAL_PASTEL),
        ("Tool",   "個別工具函式",            "search_new_books / search_courses / …",        ORANGE, ORANGE_PASTEL),
    ]
    band_h = 0.85
    band_y0 = 2.65
    band_gap = 0.18
    for i, (name, desc, ex, accent, fill) in enumerate(layers):
        y = band_y0 + i * (band_h + band_gap)
        _rounded(s, 0.85, y, 12, band_h, fill,
                 line_color=accent, line_w=2)
        _rect(s, 0.95, y + 0.10, 0.20, band_h - 0.20, accent)
        _text(s, 1.4, y, 2.5, band_h, name,
              font=FONT_TITLE, size=20, color=accent, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 4.0, y, 4.0, band_h, desc,
              font=FONT_BODY, size=14, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 8.2, y, 4.6, band_h, ex,
              font=FONT_CODE, size=12, color=MUTED, italic=True,
              anchor=MSO_ANCHOR.MIDDLE)
    page_number(s, 19, TOTAL)


def build_finale(prs):
    """Midnight bg + big punchline (final slide of main content)."""
    s = _blank_slide(prs, MIDNIGHT)
    _text(s, 0.85, 1.4, 12, 0.5, "—  一  句  話  總  結  —",
          font=FONT_BODY, size=14, color=MUTED, italic=True,
          align=PP_ALIGN.CENTER)

    _text(s, 0.85, 2.6, 12, 1.2,
          "MCP 把 function calling 標準化",
          font=FONT_TITLE, size=46, color=BG_WHITE, bold=True,
          align=PP_ALIGN.CENTER)
    _text(s, 0.85, 3.7, 12, 1.2,
          "—— LLM ↔ Tool 中間多了一個獨立 process。",
          font=FONT_TITLE, size=34, color=VIOLET, bold=True,
          align=PP_ALIGN.CENTER)

    pastel_card(s, 2.0, 5.4, 9.3, 1.2,
                accent=VIOLET, fill=RGBColor(0x1E, 0x20, 0x3A))
    _multi(s, 2.3, 5.5, 8.8, 1.05, [
        {"text": "搭配教學影片:",
         "font": FONT_BODY, "size": 14, "color": MUTED, "space_after": 4},
        {"text": "「02-mcp-connection-video.mp4 — 2:17」",
         "font": FONT_CODE, "size": 14, "color": CODE_FG, "bold": True,
         "space_after": 0},
    ])
    page_number(s, 20, TOTAL)


# ── Appendix: REST vs JSON-RPC ─────────────────────────────────────
def build_appendix_intro(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "附錄", "A P P E N D I X   ·   R E S T   v s   J S O N - R P C",
                 accent=PINK)
    slide_title(s, "為什麼 MCP 選 JSON-RPC 不選 REST?", y=0.95)
    slide_subtitle(s, "給技術 curious 老師 —— 不講也不影響理解 MCP 的核心", y=1.85)

    headers = ["面向", "REST", "JSON-RPC"]
    rows = [
        ("核心概念",  "資源(URL + HTTP 動詞)",       "動作(method 欄位)"),
        ("端點",      "每個資源一個 URL",            "單一入口,靠 method 區分"),
        ("傳輸層",    "綁定 HTTP",                  "不綁定(stdio / HTTP / WS)"),
        ("語意",      "CRUD 操作",                   "呼叫函式"),
        ("範例",      'GET /books?keyword=AI',      '{ method: "search_books" }'),
    ]

    row_h = 0.55
    table_y = 2.75
    # Header
    _rect(s, 0.85, table_y, 3.0, row_h, VIOLET_PASTEL)
    _rect(s, 3.85, table_y, 4.5, row_h, ORANGE_PASTEL)
    _rect(s, 8.35, table_y, 4.5, row_h, TEAL_PASTEL)
    _text(s, 0.85, table_y, 3.0, row_h, "面向",
          font=FONT_TITLE, size=15, color=VIOLET, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 3.85, table_y, 4.5, row_h, "REST",
          font=FONT_TITLE, size=15, color=ORANGE, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, 8.35, table_y, 4.5, row_h, "JSON-RPC",
          font=FONT_TITLE, size=15, color=TEAL, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    for i, (col1, col2, col3) in enumerate(rows):
        y = table_y + (i + 1) * row_h
        if i % 2:
            _rect(s, 0.85, y, 12, row_h, SLATE_PASTEL)
        _text(s, 1.05, y, 2.7, row_h, col1,
              font=FONT_BODY, size=13, color=INK, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 4.0, y, 4.2, row_h, col2,
              font=FONT_BODY, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, 8.5, y, 4.2, row_h, col3,
              font=FONT_BODY, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    callout_box(s, 0.85, 6.45, 12, 0.55,
                "MCP 選 JSON-RPC:傳輸層不綁 HTTP(可用 stdio),語意是「呼叫工具」而非「操作資源」",
                accent=VIOLET, fill=VIOLET_PASTEL, icon="▶", size=14)
    page_number(s, 21, TOTAL)


def build_jsonrpc_types(prs):
    s = _blank_slide(prs, BG_WHITE)
    metadata_bar(s, "附錄 · ①", "J S O N - R P C   ·   三 種 訊 息 類 型",
                 accent=PINK)
    slide_title(s, "JSON-RPC 2.0 的三種訊息類型", y=0.95)
    slide_subtitle(s, "Request / Response / Notification —— 用 id 欄位區分", y=1.85)

    items = [
        ("Request",      "Client → Server",     "帶有 id 欄位 · 期待回應",
         '{ method: "tools/call",\n  params: {...},\n  id: 1 }',          VIOLET, VIOLET_PASTEL),
        ("Response",     "Server → Client",     "帶有相同 id · 回 result 或 error",
         '{ result: {...},\n  id: 1 }',                                    TEAL,   TEAL_PASTEL),
        ("Notification", "任一方向",             "沒有 id · 不期待回應",
         '{ method:\n  "notifications/\n     initialized" }',              ORANGE, ORANGE_PASTEL),
    ]
    card_w = 3.95
    gap = 0.20
    x0 = 0.55
    for i, (name, direction, traits, example, accent, fill) in enumerate(items):
        x = x0 + i * (card_w + gap)
        _rounded(s, x, 2.7, card_w, 4.1, fill,
                 line_color=accent, line_w=2)
        _text(s, x + 0.2, 2.9, card_w - 0.4, 0.5, name,
              font=FONT_TITLE, size=18, color=accent, bold=True)
        _text(s, x + 0.2, 3.5, card_w - 0.4, 0.4, direction,
              font=FONT_BODY, size=12, color=INK, bold=True)
        _text(s, x + 0.2, 3.9, card_w - 0.4, 0.5, traits,
              font=FONT_BODY, size=12, color=INK_SOFT, italic=True)
        # Example
        _rounded(s, x + 0.2, 4.7, card_w - 0.4, 1.9, CODE_BG,
                 line_color=CODE_BG)
        _multi(s, x + 0.35, 4.85, card_w - 0.7, 1.7,
               [{"text": line, "font": FONT_CODE, "size": 11,
                 "color": CODE_STRING if i != 2 else CODE_FG,
                 "space_after": 2} for line in example.split("\n")])
    page_number(s, 22, TOTAL)


# ── Main ───────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    build_cover(prs)             # 1
    build_agenda(prs)            # 2

    build_fc_intro(prs)          # 3
    build_fc_stage12(prs)        # 4
    build_fc_stage3(prs)         # 5
    build_fc_stage4(prs)         # 6
    build_fc_stage56(prs)        # 7
    build_fc_takeaway(prs)       # 8

    build_scenario(prs)          # 9
    build_act1(prs)              # 10
    build_act2(prs)              # 11
    build_act3(prs)              # 12
    build_video_cue(prs)         # 13

    build_tools_list(prs)        # 14
    build_three_fields(prs)      # 15
    build_desc_quality(prs)      # 16

    build_client_integration(prs)  # 17
    build_llm_view(prs)            # 18
    build_architecture(prs)        # 19

    build_finale(prs)              # 20

    build_appendix_intro(prs)      # 21
    build_jsonrpc_types(prs)       # 22

    prs.save(str(PPTX))
    print(f"saved → {PPTX.name} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
