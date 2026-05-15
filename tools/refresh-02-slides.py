#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Refactor 02-how-mcp-works.pptx to align with the polished v3 video.

Plan B + appendix demotion:

  * Old slides 8-11 (lifecycle section) deleted.
  * 6 new slides inserted between old slide 2 (outline) and old slide 12:
      - 01 еҫһдёҖеҖӢжҹҘи©ўй–Ӣе§Ӣ (section divider)
      - е ҙжҷҜ: user query вҶ’ LLM tool decision
      - Act 1: Parent / Child / Sonnet badge
      - Act 2: е…©йҡҺж®өжҸЎжүӢ checkbox
      - Act 3: е·Ҙе…·е‘јеҸ« + зӯүеҖҷз¬¬ 4 еҖӢи«ӢжұӮ
      - вҶ’ еҲҮжҸӣеҲ°еҪұзүҮ
  * Old slides 3-7 (REST vs JSON-RPC) moved to end as appendix.
  * Slide 2 (outline) text refreshed to reflect new order.
  * Slide 20 (recap) vocab refreshed.

Run:
    uv run --with python-pptx python3 tools/refresh-02-slides.py

Output:
    ./02-how-mcp-works.pptx  (overwritten in place)
"""

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

REPO = Path(__file__).resolve().parent.parent
PPTX = REPO / "02-how-mcp-works.pptx"

# в”Җв”Җ Palette (Ocean Gradient + accents) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
NAVY    = RGBColor(0x1E, 0x27, 0x61)
DEEP    = RGBColor(0x06, 0x5A, 0x82)
TEAL    = RGBColor(0x1C, 0x72, 0x93)
ORANGE  = RGBColor(0xE8, 0x79, 0x3A)
BLUE    = RGBColor(0x5B, 0x8D, 0xE8)   # video v3 brighter blue
CYAN    = RGBColor(0x2D, 0xB5, 0xC9)   # video v3 brighter teal
GREEN   = RGBColor(0x5C, 0xB8, 0x5C)
RED     = RGBColor(0xD9, 0x53, 0x4F)
LIGHT_ORANGE = RGBColor(0xFC, 0xE5, 0xD1)   # faded fill for orange box + text
LIGHT_GREEN  = RGBColor(0xDB, 0xEE, 0xDB)
LIGHT_BLUE   = RGBColor(0xDD, 0xE7, 0xF9)
PAGE_BG = RGBColor(0xEF, 0xF2, 0xF7)        # light blue page background (matches kept slides)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
DARK    = RGBColor(0x1A, 0x1A, 0x1A)
MUTED   = RGBColor(0x6B, 0x6B, 0x6B)
SOFT    = RGBColor(0xE8, 0xED, 0xF3)
DIM     = RGBColor(0xC0, 0xC0, 0xC8)
CODE_BG = RGBColor(0x10, 0x1A, 0x2E)
CODE_FG = RGBColor(0xDD, 0xE4, 0xEE)

FONT_TITLE = "Arial Black"
FONT_BODY  = "Calibri"
FONT_CODE  = "Consolas"

# The existing 02-how-mcp-works.pptx is 10.0 Г— 5.625 inches.
# My build helpers were originally drafted in 13.333 Г— 7.5 coords (gen-04 style)
# so we scale all inch values down by 0.75 to fit.
S = 0.75


def _I(v):
    """Scaled Inches вҖ” maps logical (gen-04 13.333x7.5) coords to actual 10x5.625."""
    return Inches(v * S)


def _P(v):
    """Scaled Pt вҖ” fonts scale with slide too."""
    return Pt(v * S)


# в”Җв”Җ Build helpers (operate on a presentation's slides) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _add_rect(slide, x, y, w, h, color, line_color=None, line_w=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _I(x), _I(y), _I(w), _I(h)
    )
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = _P(line_w)
    sh.shadow.inherit = False
    return sh


def _add_rounded(slide, x, y, w, h, color, line_color=None, line_w=None,
                 transparency=None):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _I(x), _I(y), _I(w), _I(h)
    )
    sh.adjustments[0] = 0.12
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color
        if line_w is not None:
            sh.line.width = _P(line_w)
    sh.shadow.inherit = False
    return sh


def _add_text(slide, x, y, w, h, text, *,
              font=FONT_BODY, size=18, color=DARK,
              bold=False, italic=False, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = _I(0.05)
    tf.margin_top = tf.margin_bottom = _I(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = _P(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tb


def _add_multi(slide, x, y, w, h, paragraphs, *, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(_I(x), _I(y), _I(w), _I(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = _I(0.05)
    tf.vertical_anchor = anchor
    for i, spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if "space_after" in spec:
            p.space_after = _P(spec["space_after"])
        r = p.add_run()
        r.text = spec["text"]
        r.font.name = spec.get("font", FONT_BODY)
        r.font.size = _P(spec.get("size", 18))
        r.font.color.rgb = spec.get("color", DARK)
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
    return tb


def _blank_slide(prs, bg=WHITE):
    # Pick a blank-ish layout; fall back to layout 0 if only one exists
    layout = prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)]
    s = prs.slides.add_slide(layout)
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    sp_tree = s.shapes._spTree
    sp_tree.remove(rect._element); sp_tree.insert(2, rect._element)
    return s


def _chrome_slide(prs):
    """Standard content slide chrome: light-blue background + top orange strip.
    Matches the visual frame of the original kept slides so new + old read as
    one deck."""
    s = _blank_slide(prs, PAGE_BG)
    _add_rect(s, 0, 0, 13.333, 0.10, ORANGE)
    return s


def _apply_chrome(slide):
    """Apply the standard chrome (bg + top strip) to an existing/cleared slide."""
    _add_rect(slide, 0, 0, 13.333, 7.5, PAGE_BG)
    _add_rect(slide, 0, 0, 13.333, 0.10, ORANGE)


def _process_box(slide, x, y, w, h, header_color, header_label, body_label,
                 sonnet=False):
    """Draw a Parent/Child Process box matching the video's v3 aesthetic.

    header_label is the Chinese-or-English name shown on the colored header bar.
    body_label is the runtime name (e.g. Node.js / Python MCP Server).
    sonnet=True adds an orange LLMВ·Sonnet badge inside the box.
    """
    # Drop shadow (offset back-right)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        _I(x + 0.06), _I(y + 0.06), _I(w), _I(h)
    )
    shadow.adjustments[0] = 0.14
    shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shadow.fill.transparency = 0
    shadow.line.fill.background()
    shadow.shadow.inherit = False

    # Main box (white fill + colored border)
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, _I(x), _I(y), _I(w), _I(h)
    )
    box.adjustments[0] = 0.14
    box.fill.solid(); box.fill.fore_color.rgb = WHITE
    box.line.color.rgb = header_color
    box.line.width = _P(2.5)
    box.shadow.inherit = False

    # Header band
    hdr_h = 0.6
    hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _I(x + 0.05), _I(y + 0.05),
        _I(w - 0.10), _I(hdr_h)
    )
    hdr.fill.solid(); hdr.fill.fore_color.rgb = header_color
    hdr.line.fill.background()
    hdr.shadow.inherit = False

    # Header text
    _add_text(slide, x + 0.05, y + 0.05, w - 0.10, hdr_h,
              header_label, font=FONT_TITLE, size=20,
              color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Body label centered in remaining area; if sonnet badge will sit at the
    # bottom-right, shorten body height so the middle-anchored label doesn't
    # drop onto the badge.
    badge_zone = 0.50 if sonnet else 0
    body_y = y + 0.05 + hdr_h + 0.10
    body_h = h - hdr_h - 0.30 - badge_zone
    _add_text(slide, x + 0.05, body_y, w - 0.10, body_h,
              body_label, font=FONT_CODE, size=14,
              color=MUTED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Optional Sonnet badge at bottom-right inside box
    if sonnet:
        bw, bh = 1.5, 0.32
        bx = x + w - bw - 0.15
        by = y + h - bh - 0.15
        badge = _add_rounded(slide, bx, by, bw, bh, ORANGE,
                             line_color=ORANGE, line_w=1.5)
        _add_text(slide, bx, by, bw, bh, "LLM В· Sonnet",
                  font=FONT_CODE, size=11, color=WHITE, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _arrow(slide, x1, y1, x2, y2, color, line_w=3):
    """Straight arrow from (x1,y1) to (x2,y2), inches."""
    conn = slide.shapes.add_connector(
        1,  # straight
        _I(x1), _I(y1), _I(x2), _I(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = _P(line_w)
    # Add arrowhead via XML
    from pptx.oxml.ns import qn
    ln = conn.line._get_or_add_ln()
    # Remove any existing tailEnd / headEnd
    for tag in ("a:tailEnd", "a:headEnd"):
        for el in ln.findall(qn(tag)):
            ln.remove(el)
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    tail = etree.SubElement(ln, qn("a:tailEnd"), nsmap=nsmap)
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return conn


def _pipe_pair(slide, x1, y1, x2, y2, color=MUTED, gap=0.18):
    """Two parallel lines representing bidirectional pipes."""
    for dy in (-gap, +gap):
        conn = slide.shapes.add_connector(
            1, _I(x1), _I(y1 + dy), _I(x2), _I(y2 + dy)
        )
        conn.line.color.rgb = color
        conn.line.width = _P(1.5)


# в”Җв”Җ New slide content в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def build_slide_section_divider(prs, number, title, subtitle):
    """Section divider вҖ” compact style (matches original kept dividers).
    Teal bg + thin orange top strip + small number top-left + title below."""
    s = _blank_slide(prs, TEAL)
    _add_rect(s, 0, 0, 13.333, 0.10, ORANGE)
    # Number (orange, small upper-left)
    _add_text(s, 0.85, 2.0, 4.0, 1.0, number,
              font=FONT_TITLE, size=46, color=ORANGE, bold=True)
    # Title (white, large, below number; full width avoids overflow)
    _add_text(s, 0.85, 2.95, 12, 1.0, title,
              font=FONT_TITLE, size=42, color=WHITE, bold=True)
    # Subtitle (soft, below title)
    _add_text(s, 0.85, 4.05, 12, 0.7, subtitle,
              font=FONT_BODY, size=20, color=SOFT)
    return s


def build_slide_scenario(prs):
    """Scenario: user query вҶ’ LLM tool decision."""
    s = _chrome_slide(prs)
    # Title
    _add_text(s, 0.75, 0.4, 12, 0.7, "е ҙжҷҜ:дҪҝз”ЁиҖ…е•ҸдәҶдёҖеҖӢе•ҸйЎҢ,LLM жұәе®ҡиҰҒе‘јеҸ«е·Ҙе…·",
              font=FONT_TITLE, size=26, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    # User chat bubble (top-right area)
    bub_x, bub_y = 5.5, 1.8
    bub_w, bub_h = 7.0, 1.1
    _add_text(s, bub_x + bub_w - 1.5, bub_y - 0.45, 1.5, 0.35, "дҪҝз”ЁиҖ…",
              font=FONT_BODY, size=14, color=MUTED, align=PP_ALIGN.RIGHT)
    _add_rounded(s, bub_x, bub_y, bub_w, bub_h, BLUE,
                 line_color=BLUE, line_w=2.5)
    _add_text(s, bub_x, bub_y, bub_w, bub_h,
              "е№«жҲ‘жҹҘжңҖиҝ‘ең–жӣёйӨЁжңүд»Җйәјж–°жӣё",
              font=FONT_BODY, size=22, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Diagonal arrow user вҶ’ LLM
    _arrow(s, 6.5, 3.0, 4.5, 4.2, MUTED, line_w=2)

    # LLM thinking box (bottom-left) вҖ” faded orange fill so orange text stays visible
    llm_x, llm_y = 0.8, 4.3
    llm_w, llm_h = 7.5, 1.8
    _add_rounded(s, llm_x, llm_y, llm_w, llm_h, LIGHT_ORANGE,
                 line_color=ORANGE, line_w=2.5)
    # Fade orange fill (use line only with white fill)
    _add_multi(s, llm_x + 0.3, llm_y + 0.25, llm_w - 0.6, llm_h - 0.5, [
        {"text": "е—ҜвҖҰйңҖиҰҒжҹҘж–°жӣёиіҮж–ҷеә«",
         "font": FONT_BODY, "size": 20, "color": DARK},
        {"text": "вҶ’ е‘јеҸ« search_new_books е·Ҙе…·",
         "font": FONT_CODE, "size": 20, "color": ORANGE, "bold": True,
         "space_after": 0},
    ])
    _add_text(s, llm_x, llm_y + llm_h + 0.05, 4, 0.35, "LLM (Sonnet)",
              font=FONT_CODE, size=14, color=MUTED)

    # Bottom hook
    _add_text(s, 0.75, 6.55, 12.5, 0.5,
              "дҪҶ LLM иҮӘе·ұдёҚжңғе‘јеҸ« вҖ” иҰҒйқ  Parent еҺ»и·ҹ Child и¬ӣи©ұ",
              font=FONT_BODY, size=20, color=RED, bold=True, italic=True,
              align=PP_ALIGN.CENTER)
    return s


def build_slide_act1(prs):
    """Act 1: Parent / Child boxes + LLM Sonnet badge + pipes."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7, "Act 1 вҖ” Parent й–ӢеҮә Child",
              font=FONT_TITLE, size=28, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    # Spawn cue at top
    _add_text(s, 4.5, 1.45, 6, 0.5, "в–¶ й–ӢеҮә Child",
              font=FONT_CODE, size=22, color=ORANGE, bold=True,
              align=PP_ALIGN.CENTER)

    # Parent box (left)
    _process_box(s, 0.8, 2.1, 4.2, 2.5, BLUE,
                 "Parent Process", "Node.js", sonnet=True)
    # Child box (right)
    _process_box(s, 8.3, 2.1, 4.2, 2.5, CYAN,
                 "Child Process", "Python MCP Server", sonnet=False)

    # Two pipes between them
    _pipe_pair(s, 5.05, 3.35, 8.25, 3.35, MUTED, gap=0.20)

    # Bottom bullets
    _add_multi(s, 0.75, 5.0, 12.5, 2.2, [
        {"text": "вҖў Parent (Node.js) з”Ё spawn() й–ӢеҮә Child (Python) еӯҗзЁӢеәҸ",
         "font": FONT_BODY, "size": 20, "color": DARK,
         "space_after": 8},
        {"text": "вҖў LLM (Sonnet) жҳҜ Parent иЈЎзҡ„дёҖеҖӢ client component вҖ” е®ғдҪҸеңЁ Parent е…§,дёҚжңғиҮӘе·ұеҮәеҺ»е‘јеҸ«е·Ҙе…·",
         "font": FONT_BODY, "size": 20, "color": DARK,
         "space_after": 8},
        {"text": "вҖў е…©жўқ stdio pipe дёІйҖҡ вҶ’ Parent вҮ„ Child еҸҜд»Ҙйӣҷеҗ‘и¬ӣи©ұ",
         "font": FONT_BODY, "size": 20, "color": DARK,
         "space_after": 0},
    ])
    return s


def build_slide_act2(prs):
    """Act 2: е…©йҡҺж®өжҸЎжүӢ with checkbox visualization."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7, "Act 2 вҖ” е…©йҡҺж®өжҸЎжүӢ",
              font=FONT_TITLE, size=28, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    # Left side: Parent вҶ’ Child message exchange
    _process_box(s, 0.8, 2.0, 4.0, 1.5, BLUE, "Parent", "Node.js",
                 sonnet=False)
    _process_box(s, 0.8, 4.3, 4.0, 1.5, CYAN, "Child", "Python",
                 sonnet=False)

    # Two arrows: Parent вҶ’ Child (orange) and Child вҶ’ Parent (green)
    _arrow(s, 2.8, 3.5, 2.8, 4.3, ORANGE, line_w=3)
    _arrow(s, 2.8, 5.8, 2.8, 6.5, GREEN, line_w=3)

    _add_text(s, 5.0, 3.65, 4, 0.4, "вҶ’ жҲ‘дҫҶдәҶ / зөҰжҲ‘е·Ҙе…·жё…е–®",
              font=FONT_BODY, size=15, color=ORANGE, bold=True)
    _add_text(s, 5.0, 5.85, 4, 0.4, "вҶҗ жҲ‘иғҪеҒҡд»Җйәј / 8 еҖӢе·Ҙе…·",
              font=FONT_BODY, size=15, color=GREEN, bold=True)

    # Right side: Two-step checklist
    panel_x, panel_y, panel_w, panel_h = 8.0, 2.0, 4.7, 3.8
    _add_rounded(s, panel_x, panel_y, panel_w, panel_h, SOFT,
                 line_color=MUTED, line_w=1.5)
    _add_text(s, panel_x + 0.3, panel_y + 0.2, panel_w - 0.6, 0.5,
              "жҸЎжүӢйҖІеәҰ", font=FONT_TITLE, size=18, color=NAVY, bold=True)

    # Step 1: checked green
    sx = panel_x + 0.4
    sy1 = panel_y + 0.9
    _add_rounded(s, sx, sy1, 0.5, 0.5, GREEN, line_color=GREEN, line_w=2)
    _add_text(s, sx, sy1, 0.5, 0.5, "вң“",
              font=FONT_TITLE, size=22, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(s, sx + 0.7, sy1, 3.5, 0.5, "Step 1:жү“йҒҺжӢӣе‘ј",
              font=FONT_BODY, size=20, color=GREEN, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)

    # Step 2: checked green
    sy2 = sy1 + 0.95
    _add_rounded(s, sx, sy2, 0.5, 0.5, GREEN, line_color=GREEN, line_w=2)
    _add_text(s, sx, sy2, 0.5, 0.5, "вң“",
              font=FONT_TITLE, size=22, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text(s, sx + 0.7, sy2, 3.5, 0.5, "Step 2:жӢҝеҲ°е·Ҙе…·жё…е–®",
              font=FONT_BODY, size=20, color=GREEN, bold=True,
              anchor=MSO_ANCHOR.MIDDLE)

    # Bottom note inside panel
    _add_text(s, panel_x + 0.3, panel_y + panel_h - 1.0, panel_w - 0.6, 0.6,
              "е…©еҖӢйғҪжү“еӢҫ жүҚз®—зңҹзҡ„йҖЈдёҠ",
              font=FONT_BODY, size=18, color=NAVY, bold=True, italic=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom takeaway
    _add_text(s, 0.75, 6.4, 12.5, 0.5,
              "йҮҚй»һ:дёӯй–“жҷӮеҲ»(Step 1 вң“ дҪҶ Step 2 вң—)Parent йӮ„дёҚзҹҘйҒ“е·Ҙе…·жё…е–®,йҖҒ tool call д№ҹжІ’з”Ё",
              font=FONT_BODY, size=18, color=MUTED, italic=True,
              align=PP_ALIGN.CENTER)
    return s


def build_slide_act3(prs):
    """Act 3: tool call + waiting list + result back."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7, "Act 3 вҖ” е·Ҙе…·е‘јеҸ«",
              font=FONT_TITLE, size=28, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    # Parent + Child boxes
    _process_box(s, 0.8, 1.7, 3.6, 1.4, BLUE, "Parent", "Node.js",
                 sonnet=True)
    _process_box(s, 9.0, 1.7, 3.6, 1.4, CYAN, "Child", "Python",
                 sonnet=False)

    # Arrow Parent вҶ’ Child (orange) with token
    _arrow(s, 4.4, 2.4, 9.0, 2.4, ORANGE, line_w=3)
    _add_rounded(s, 5.5, 2.05, 2.5, 0.55, ORANGE,
                 line_color=ORANGE, line_w=2)
    _add_text(s, 5.5, 2.05, 2.5, 0.55,
              'search_new_books("ж–°жӣё")',
              font=FONT_CODE, size=11, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Waiting list box below Parent вҖ” light orange so orange text reads
    _add_rounded(s, 0.8, 3.3, 3.6, 0.6, LIGHT_ORANGE,
                 line_color=ORANGE, line_w=2)
    _add_text(s, 0.8, 3.3, 3.6, 0.6,
              "зӯүеҖҷз¬¬ 4 еҖӢи«ӢжұӮ В· 30s е…§",
              font=FONT_BODY, size=14, color=ORANGE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrow Child вҶ’ Parent (green) with result token
    _arrow(s, 9.0, 4.5, 4.4, 4.5, GREEN, line_w=3)
    _add_rounded(s, 5.5, 4.18, 2.5, 0.55, GREEN, line_color=GREEN, line_w=2)
    _add_text(s, 5.5, 4.18, 2.5, 0.55,
              "зөҗжһң В· 10 зӯҶж–°жӣё",
              font=FONT_CODE, size=12, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Done waiting list (resolved) вҖ” light green so green text reads
    _add_rounded(s, 0.8, 5.0, 3.6, 0.6, LIGHT_GREEN,
                 line_color=GREEN, line_w=2)
    _add_text(s, 0.8, 5.0, 3.6, 0.6,
              "з¬¬ 4 еҖӢи«ӢжұӮ вң“ е®ҢжҲҗ",
              font=FONT_BODY, size=14, color=GREEN, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom flow recap
    _add_text(s, 0.75, 5.95, 12.5, 0.5,
              "е®Ңж•ҙжөҒзЁӢ:LLM жұәе®ҡ вҶ’ token йЈӣеҺ» Child вҶ’ Child еҹ·иЎҢ вҶ’ "
              "зөҗжһңйЈӣеӣһ вҶ’ Parent жүҫеҲ°е°ҚжҮүи«ӢжұӮ вҶ’ еӣһзөҰдҪҝз”ЁиҖ…",
              font=FONT_BODY, size=16, color=DARK,
              align=PP_ALIGN.CENTER)

    _add_text(s, 0.75, 6.55, 12.5, 0.5,
              "гҖҢз¬¬ 4 еҖӢи«ӢжұӮгҖҚжҳҜ JSON-RPC зөҰжҜҸеҖӢ request зҡ„з·Ёиҷҹ,и®“ response иғҪе°Қеҫ—еӣһдҫҶ",
              font=FONT_BODY, size=14, color=MUTED, italic=True,
              align=PP_ALIGN.CENTER)
    return s


def build_slide_video_cue(prs):
    """Cue slide: switch to the polished video."""
    s = _blank_slide(prs, NAVY)
    _add_text(s, 0.75, 1.5, 12, 1.2, "вҶ’  еҲҮжҸӣеҲ°еҪұзүҮ",
              font=FONT_TITLE, size=64, color=ORANGE, bold=True)
    _add_text(s, 0.75, 3.0, 12, 0.6,
              "02-mcp-connection-video.mp4   В·   2:17",
              font=FONT_CODE, size=22, color=WHITE)
    _add_rect(s, 0.75, 3.85, 2.5, 0.04, ORANGE)

    _add_multi(s, 0.75, 4.0, 12, 3.0, [
        {"text": "еҪұзүҮйҮҚй»һ:",
         "font": FONT_BODY, "size": 22, "color": WHITE, "bold": True,
         "space_after": 14},
        {"text": "1. е ҙжҷҜй–Ӣе ҙ вҖ” дҪҝз”ЁиҖ…е•ҸгҖҢж–°жӣёгҖҚ вҶ’ LLM еҲӨж–·иҰҒе‘јеҸ«е·Ҙе…·",
         "font": FONT_BODY, "size": 19, "color": SOFT, "space_after": 8},
        {"text": "2. Act 1 вҖ” Parent / Child зҡ„ spawn иҲҮ pipe;LLM (Sonnet) дҪҸеңЁ Parent е…§",
         "font": FONT_BODY, "size": 19, "color": SOFT, "space_after": 8},
        {"text": "3. Act 2 вҖ” е…©йҡҺж®өжҸЎжүӢ:checkbox вҳҗжү“йҒҺжӢӣе‘ј вҶ’ вң“жү“йҒҺжӢӣе‘ј вҶ’ вҳҗжӢҝе·Ҙе…· вҶ’ вң“жӢҝе·Ҙе…·",
         "font": FONT_BODY, "size": 19, "color": SOFT, "space_after": 8},
        {"text": "4. Act 3 вҖ” е·Ҙе…·е‘јеҸ«:token йЈӣеҺ»гҖҒзӯүеҖҷз¬¬ 4 еҖӢи«ӢжұӮгҖҒзөҗжһңйЈӣеӣһ",
         "font": FONT_BODY, "size": 19, "color": SOFT, "space_after": 8},
        {"text": "5. Frame closure вҖ” гҖҢжҹҘж–°жӣёгҖҚвҶ’ Child еҹ·иЎҢ вҶ’ 10 зӯҶж–°жӣё вҶ’ дҪҝз”ЁиҖ…",
         "font": FONT_BODY, "size": 19, "color": SOFT, "space_after": 0},
    ])
    return s


def build_appendix_divider(prs):
    """Appendix divider вҖ” same compact style as other section dividers."""
    s = _blank_slide(prs, TEAL)
    _add_rect(s, 0, 0, 13.333, 0.10, ORANGE)
    _add_text(s, 0.85, 2.0, 4.0, 1.0, "йҷ„йҢ„",
              font=FONT_TITLE, size=46, color=ORANGE, bold=True)
    _add_text(s, 0.85, 2.95, 12, 1.0, "Appendix",
              font=FONT_TITLE, size=42, color=WHITE, bold=True)
    _add_multi(s, 0.85, 4.05, 12, 2.2, [
        {"text": "зөҰжҠҖиЎ“ curious иҖҒеё«зҡ„зҙ°зҜҖ",
         "font": FONT_BODY, "size": 20, "color": SOFT, "space_after": 10},
        {"text": "вҖў REST vs JSON-RPC жҜ”ијғ",
         "font": FONT_BODY, "size": 17, "color": SOFT, "space_after": 4},
        {"text": "вҖў JSON-RPC дёүзЁ®иЁҠжҒҜйЎһеһӢ (Request / Response / Notification)",
         "font": FONT_BODY, "size": 17, "color": SOFT, "space_after": 0},
    ])
    return s


# в”Җв”Җ Slide list manipulation (XML level) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _slides_xml(prs):
    return prs.slides._sldIdLst


def _slide_ids(prs):
    return list(_slides_xml(prs))


def _move_slide(prs, from_index, to_index):
    """Move slide currently at from_index to to_index."""
    xml = _slides_xml(prs)
    slides = list(xml)
    s = slides[from_index]
    xml.remove(s)
    # After removal, indices shift
    target = to_index if to_index < from_index else to_index - 1
    if target >= len(xml):
        xml.append(s)
    else:
        xml.insert(target, s)


def _delete_slide(prs, index):
    xml = _slides_xml(prs)
    slides = list(xml)
    xml.remove(slides[index])


# в”Җв”Җ Slide 2 outline rewrite в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _clear_slide_shapes(slide):
    """Remove all shapes from a slide (used to rebuild slide 2 / 20)."""
    sp_tree = slide.shapes._spTree
    # Keep first 2 elements (nvGrpSpPr + grpSpPr) which are required
    children = list(sp_tree)
    for el in children:
        tag = el.tag.split("}")[-1]
        if tag in ("sp", "pic", "graphicFrame", "cxnSp", "grpSp"):
            sp_tree.remove(el)


def rebuild_slide_2(slide):
    """Rebuild slide 2 outline to reflect new structure."""
    _clear_slide_shapes(slide)
    _apply_chrome(slide)
    _add_text(slide, 0.75, 0.5, 12, 0.8, "жң¬и¬ӣеӨ§з¶ұ",
              font=FONT_TITLE, size=36, color=NAVY, bold=True)
    _add_rect(slide, 0.75, 1.40, 3, 0.04, ORANGE)

    items = [
        ("01", "Function Calling жҖҺйәјйҒӢдҪң",
         "LLM еҸӘеҗҗеӯ—дёІ,зңҹзҡ„еҹ·иЎҢзҡ„жҳҜеӨ–йқўзҡ„ harness"),
        ("02", "еҫһдёҖеҖӢжҹҘи©ўй–Ӣе§Ӣ",
         "е ҙжҷҜ вҶ’ Parent й–ӢеҮә Child вҶ’ е…©йҡҺж®өжҸЎжүӢ вҶ’ е·Ҙе…·е‘јеҸ«"),
        ("03", "Tool иЁ»еҶҠиҲҮжҸҸиҝ°",
         "JSON Schema В· description еҰӮдҪ•еҪұйҹҝ LLM йҒёж“Ү"),
        ("04", "Client ж•ҙеҗҲж©ҹеҲ¶",
         "е·Ҙе…·жё…е–®еҰӮдҪ•жіЁе…Ҙ LLM зҡ„ system prompt"),
        ("йҷ„йҢ„", "жҠҖиЎ“зҙ°зҜҖ(жҷӮй–“еӨ е°ұи¬ӣ,дёҚеӨ е°ұи·і)",
         "REST vs JSON-RPC В· JSON-RPC дёүзЁ®иЁҠжҒҜйЎһеһӢ"),
    ]
    y = 1.7
    for num, title, sub in items:
        _add_text(slide, 0.8, y, 1.7, 0.7, num,
                  font=FONT_TITLE, size=36, color=ORANGE, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, 2.7, y, 10, 0.55, title,
                  font=FONT_TITLE, size=23, color=NAVY, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, 2.7, y + 0.55, 10, 0.45, sub,
                  font=FONT_BODY, size=15, color=MUTED,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += 1.05


def _code_block(slide, x, y, w, h, lines, size=11, padding=0.18):
    """Dark code block; `lines` is list[str] or list[(text, color)]."""
    _add_rect(slide, x, y, w, h, CODE_BG)
    paragraphs = []
    for line in lines:
        if isinstance(line, tuple):
            text, color = line
        else:
            text, color = line, CODE_FG
        paragraphs.append({
            "text": text if text else " ",
            "font": FONT_CODE,
            "size": size,
            "color": color,
            "space_after": 1,
        })
    _add_multi(slide, x + 0.25, y + padding, w - 0.45, h - padding * 2,
               paragraphs)


# в”Җв”Җ Section 01 (new): Function Calling жҖҺйәјйҒӢдҪң в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def build_fc_setup(prs):
    """Setup slide: one concrete example + the core punchline."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7,
              "з”ЁдёҖеҖӢе…·й«”дҫӢеӯҗжҗһжҮӮ",
              font=FONT_TITLE, size=28, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    # User query bubble
    bub_x, bub_y, bub_w, bub_h = 2.5, 1.9, 8.0, 1.0
    _add_text(s, bub_x, bub_y - 0.45, bub_w, 0.35, "дҪҝз”ЁиҖ…ијёе…Ҙ",
              font=FONT_BODY, size=14, color=MUTED,
              align=PP_ALIGN.CENTER)
    _add_rounded(s, bub_x, bub_y, bub_w, bub_h, BLUE,
                 line_color=BLUE, line_w=2.5)
    _add_text(s, bub_x, bub_y, bub_w, bub_h,
              "жҲ‘е®¶зӣ®йҢ„еә•дёӢжңүе№ҫеҖӢ .py жӘ”пјҹ",
              font=FONT_BODY, size=24, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Big punchline
    _add_text(s, 0.75, 3.5, 12, 0.9,
              "LLM еҫһй ӯеҲ°е°ҫеҸӘеҒҡдёҖд»¶дәӢ вҖ”вҖ” еҗҗеӯ—дёІ",
              font=FONT_TITLE, size=32, color=ORANGE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_multi(s, 0.75, 4.9, 12, 2.0, [
        {"text": "жүҖжңүгҖҢзңҹзҡ„еҹ·иЎҢгҖҚзҡ„еӢ•дҪң вҖ”вҖ” и·‘ findгҖҒзҷј HTTP requestгҖҒж”№жӘ”гҖҒй–ӢзҖҸиҰҪеҷЁ вҖ”вҖ” йғҪзҷјз”ҹеңЁ LLM еӨ–йқўзҡ„дёҖж®өзЁӢејҸиЈЎгҖӮ",
         "font": FONT_BODY, "size": 18, "color": DARK,
         "align": PP_ALIGN.CENTER, "space_after": 14},
        {"text": "жҲ‘еҖ‘еҸ«е®ғ harness  (= agent runtime / agent loop runner)",
         "font": FONT_BODY, "size": 20, "color": NAVY, "bold": True,
         "align": PP_ALIGN.CENTER, "space_after": 0},
    ])
    return s


def build_fc_stage12(prs):
    """Stage 1+2: developer prepares schema + sends API call."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7,
              "Stage 1+2 вҖ” й–ӢзҷјиҖ…жә–еӮҷе·Ҙе…·жҸҸиҝ°,йҖҒеҮә API call",
              font=FONT_TITLE, size=24, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    code_lines = [
        ("# й–ӢзҷјиҖ…еҜ«зҡ„ Python", MUTED),
        ("tools = [{",                                                CODE_FG),
        ('  "name": "execute_bash",',                                 GREEN),
        ('  "description": "Execute a bash command. Returns stdout.",', GREEN),
        ('  "input_schema": {',                                       CODE_FG),
        ('    "type": "object",',                                     CODE_FG),
        ('    "properties": {',                                       CODE_FG),
        ('      "command": {"type": "string"}',                       CODE_FG),
        ("    },",                                                    CODE_FG),
        ('    "required": ["command"]',                               CODE_FG),
        ("  }",                                                       CODE_FG),
        ("}]",                                                        CODE_FG),
        ("",                                                          CODE_FG),
        ('messages = [{"role": "user",',                              CODE_FG),
        ('             "content": "жҲ‘е®¶зӣ®йҢ„еә•дёӢжңүе№ҫеҖӢ .py жӘ”пјҹ"}]',     CODE_FG),
        ("",                                                          CODE_FG),
        ("response = client.messages.create(",                        CODE_FG),
        ('    model="claude-opus-4-7",',                              CODE_FG),
        ("    tools=tools,        вҶҗ жҠҠе·Ҙе…· schema дёҖиө·йҖҒ",             ORANGE),
        ("    messages=messages,",                                    CODE_FG),
        (")",                                                         CODE_FG),
    ]
    _code_block(s, 0.75, 1.55, 7.5, 5.3, code_lines)

    _add_multi(s, 8.6, 1.7, 4.5, 5.2, [
        {"text": "tools еҸӘжҳҜдёҖеҖӢ dict",
         "font": FONT_TITLE, "size": 17, "color": NAVY, "bold": True,
         "space_after": 4},
        {"text": "SDK жҠҠе®ғеәҸеҲ—еҢ–жҲҗ JSON,еЎһйҖІ API зҡ„ request bodyгҖӮ",
         "font": FONT_BODY, "size": 14, "color": DARK,
         "space_after": 16},

        {"text": "LLM зңӢеҲ°зҡ„жҳҜйҖҷж®ө JSON",
         "font": FONT_TITLE, "size": 17, "color": NAVY, "bold": True,
         "space_after": 4},
        {"text": "зҝ»жҲҗ prompt еҫҢ LLM и®ҖеҲ°гҖҢжңүдёҖеҖӢеҸ« execute_bash зҡ„жқұиҘҝ,жҸҸиҝ°жҳҜйҖҷжЁЈ,еҸғж•ёй•·йҖҷжЁЈгҖҚгҖӮ",
         "font": FONT_BODY, "size": 14, "color": DARK,
         "space_after": 16},

        {"text": "вң— жІ’жңүз¶Ғе®ҡзңҹеҮҪејҸжҢҮжЁҷ",
         "font": FONT_TITLE, "size": 17, "color": RED, "bold": True,
         "space_after": 4},
        {"text": "LLM жӢҝдёҚеҲ° Python еҮҪејҸ reference,д№ҹзў°дёҚеҲ° OSгҖӮе®ғе°ұжҳҜи®ҖдәҶдёҖж®өж–Үеӯ—гҖӮ",
         "font": FONT_BODY, "size": 14, "color": DARK,
         "space_after": 0},
    ])
    return s


def build_fc_stage3(prs):
    """Stage 3: LLM returns JSON вҖ” the KEY teaching moment."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7,
              "Stage 3 вҖ” LLM еӣһдәҶдёҖж®ө JSON",
              font=FONT_TITLE, size=26, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    json_lines = [
        ("{",                                                         CODE_FG),
        ('  "id": "msg_abc123",',                                     CODE_FG),
        ('  "stop_reason": "tool_use",        вҶҗ ж——иҷҹ:зөҰ harness зңӢ',  ORANGE),
        ('  "content": [',                                            CODE_FG),
        ('    {',                                                     CODE_FG),
        ('      "type": "text",',                                     CODE_FG),
        ('      "text": "и®“жҲ‘з”Ё find е№«дҪ ж•ёдёҖдёӢгҖӮ"',                   GREEN),
        ('    },',                                                    CODE_FG),
        ('    {',                                                     CODE_FG),
        ('      "type": "tool_use",',                                 CODE_FG),
        ('      "id": "toolu_01XYZ",',                                CODE_FG),
        ('      "name": "execute_bash",',                             CODE_FG),
        ('      "input": {',                                          CODE_FG),
        ('        "command": "find ~ -name \'*.py\' -type f | wc -l"', GREEN),
        ('      }',                                                   CODE_FG),
        ('    }',                                                     CODE_FG),
        ('  ]',                                                       CODE_FG),
        ('}',                                                         CODE_FG),
    ]
    _code_block(s, 0.75, 1.55, 8.6, 4.3, json_lines)

    _add_multi(s, 9.6, 1.7, 3.6, 4.2, [
        {"text": "stop_reason: tool_use",
         "font": FONT_CODE, "size": 14, "color": ORANGE, "bold": True,
         "space_after": 4},
        {"text": "гҖҢжҲ‘еҒңдёӢдҫҶдәҶ,еӣ зӮәз”ўеҮәдәҶдёҖеҖӢ tool_use,дҪ жҺҘжүӢгҖӮгҖҚ",
         "font": FONT_BODY, "size": 13, "color": MUTED, "italic": True,
         "space_after": 16},

        {"text": "command йҖҷдёІеӯ—",
         "font": FONT_TITLE, "size": 14, "color": GREEN, "bold": True,
         "space_after": 4},
        {"text": "жҳҜ LLM token-by-token з”ҹеҮәдҫҶзҡ„;constrained decoding зўәдҝқ JSON дёҖе®ҡеҗҲжі•гҖӮ",
         "font": FONT_BODY, "size": 13, "color": DARK,
         "space_after": 0},
    ])

    # Big punchline at bottom
    _add_rounded(s, 0.75, 6.05, 12, 0.85, LIGHT_ORANGE,
                 line_color=ORANGE, line_w=2)
    _add_text(s, 0.75, 6.05, 12, 0.85,
              "LLM жІ’жңүгҖҢе‘јеҸ«гҖҚд»»дҪ•жқұиҘҝ вҖ”вҖ” е®ғе°ұжҳҜеҗҗдәҶдёҖж®ө JSON еӯ—дёІгҖӮ",
              font=FONT_TITLE, size=20, color=ORANGE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return s


def build_fc_stage4(prs):
    """Stage 4: Harness actually executes вҖ” the 'eval LLM output' moment."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7,
              "Stage 4 вҖ” Harness зңҹзҡ„еҹ·иЎҢ",
              font=FONT_TITLE, size=26, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    code_lines = [
        ("# Harness жӢҝеҲ° LLM зҡ„ tool_use еҫҢ,и·‘ dispatcher", MUTED),
        ("for block in response.content:",            CODE_FG),
        ('    if block.type == "tool_use":',          CODE_FG),
        ('        cmd = block.input["command"]',      CODE_FG),
        ("",                                          CODE_FG),
        ("        # йҖҷдёҖиЎҢжүҚжҳҜгҖҢзңҹзҡ„дёӢеҺ»еҹ·иЎҢгҖҚ",       MUTED),
        ("        result = subprocess.run(",          ORANGE),
        ("            cmd,",                          CODE_FG),
        ("            shell=True,         вҶҗ LLM еӯ—дёІ вҶ’ /bin/bash", ORANGE),
        ("            capture_output=True,",          CODE_FG),
        ("            text=True,",                    CODE_FG),
        ("            timeout=30,",                   CODE_FG),
        ("        )",                                 CODE_FG),
        ('        tool_output = result.stdout.strip()    # "42"', CODE_FG),
    ]
    _code_block(s, 0.75, 1.55, 8.2, 3.7, code_lines)

    # Right side: door metaphor
    _add_multi(s, 9.2, 1.7, 4.0, 3.5, [
        {"text": "гҖҢж–Үеӯ—е®Үе®ҷгҖҚ",
         "font": FONT_TITLE, "size": 16, "color": BLUE, "bold": True,
         "align": PP_ALIGN.CENTER, "space_after": 2},
        {"text": "= LLM еҗҗзҡ„еӯ—дёІ",
         "font": FONT_BODY, "size": 13, "color": MUTED,
         "align": PP_ALIGN.CENTER, "space_after": 14},

        {"text": "вҶ“  йҖҷйҒ“й–Җ  вҶ“",
         "font": FONT_TITLE, "size": 20, "color": ORANGE, "bold": True,
         "align": PP_ALIGN.CENTER, "space_after": 4},
        {"text": "subprocess.run(...)",
         "font": FONT_CODE, "size": 13, "color": ORANGE, "bold": True,
         "align": PP_ALIGN.CENTER, "space_after": 14},

        {"text": "гҖҢзңҹеҜҰйӣ»и…ҰгҖҚ",
         "font": FONT_TITLE, "size": 16, "color": NAVY, "bold": True,
         "align": PP_ALIGN.CENTER, "space_after": 2},
        {"text": "= shell fork + exec",
         "font": FONT_BODY, "size": 13, "color": MUTED,
         "align": PP_ALIGN.CENTER, "space_after": 0},
    ])

    # Bottom warning box
    _add_rounded(s, 0.75, 5.45, 12, 1.5,
                 RGBColor(0xFD, 0xE5, 0xE5),
                 line_color=RED, line_w=2)
    _add_multi(s, 1.0, 5.6, 11.5, 1.3, [
        {"text": "вҡ   е®үе…ЁиӯҰе‘Ҡ вҖ” жң¬иіӘдёҠжҳҜ eval(LLM ијёеҮә)",
         "font": FONT_TITLE, "size": 16, "color": RED, "bold": True,
         "space_after": 6},
        {"text": "иӢҘ prompt injection ж”»ж“Ҡ,LLM жңғз”ҹеҮә rm -rf ~,harness дёҚй•·зңје°ұз…§и·‘гҖӮ",
         "font": FONT_BODY, "size": 14, "color": DARK,
         "space_after": 4},
        {"text": "Production зі»зөұеҝ…й ҲеҠ  sandbox / docker / allow-list / human-in-the-loop confirmationгҖӮ",
         "font": FONT_BODY, "size": 14, "color": DARK, "italic": True,
         "space_after": 0},
    ])
    return s


def build_fc_stage56(prs):
    """Stage 5+6: tool_result fed back, LLM produces final answer."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7,
              "Stage 5+6 вҖ” зөҗжһңеЎһеӣһ,LLM зөҰжңҖзөӮзӯ”жЎҲ",
              font=FONT_TITLE, size=24, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    _add_text(s, 0.75, 1.55, 8, 0.4, "messages йҷЈеҲ—зҸҫеңЁй•·йҖҷжЁЈ:",
              font=FONT_TITLE, size=14, color=NAVY, bold=True)

    msg_lines = [
        ("[user]",                                                BLUE),
        ('    "жҲ‘е®¶зӣ®йҢ„еә•дёӢжңүе№ҫеҖӢ .py жӘ”пјҹ"',                       CODE_FG),
        ("",                                                       CODE_FG),
        ("[assistant]",                                            ORANGE),
        ('    "и®“жҲ‘з”Ё find е№«дҪ ж•ёдёҖдёӢгҖӮ"',                          CODE_FG),
        ('    + tool_use { execute_bash, "find ~ ..." }',          CODE_FG),
        ("",                                                       CODE_FG),
        ("[user]   вҶҗ жіЁж„Ҹ:role жҳҜ user!",                         GREEN),
        ('    tool_result { id=toolu_01XYZ, content="42" }',       CODE_FG),
        ("",                                                       CODE_FG),
        ("[assistant]   stop_reason: end_turn",                    ORANGE),
        ('    "дҪ е®¶зӣ®йҢ„еә•дёӢзёҪе…ұжңү 42 еҖӢ .py жӘ”гҖӮ"',                  CODE_FG),
    ]
    _code_block(s, 0.75, 2.0, 7.8, 4.4, msg_lines)

    _add_multi(s, 8.9, 2.0, 4.2, 4.4, [
        {"text": "tool_result з”Ё role: user йҖҒеӣһ",
         "font": FONT_TITLE, "size": 15, "color": GREEN, "bold": True,
         "space_after": 4},
        {"text": "еҫһ model зҡ„и§’еәҰ,е®ғдёӯж–·иҮӘе·ұеҺ»е•ҸдәҶеӨ–йқўзҡ„дё–з•Ң,еӨ–йқўзҡ„дё–з•ҢеӣһдәҶдёҖеҸҘ,зҸҫеңЁијӘеҲ° model з№јзәҢи¬ӣгҖӮ",
         "font": FONT_BODY, "size": 13, "color": DARK,
         "space_after": 18},

        {"text": "stop_reason: end_turn",
         "font": FONT_CODE, "size": 14, "color": ORANGE, "bold": True,
         "space_after": 4},
        {"text": "harness зңӢеҲ°йҖҷеҖӢ(дёҚжҳҜ tool_use)е°ұзҹҘйҒ“:иҝҙеңҲзөҗжқҹдәҶгҖӮ",
         "font": FONT_BODY, "size": 13, "color": DARK,
         "space_after": 18},

        {"text": "е®Ңж•ҙ loop",
         "font": FONT_TITLE, "size": 14, "color": NAVY, "bold": True,
         "space_after": 4},
        {"text": "tool_use вҶ’ run вҶ’ tool_result вҶ’ вҖҰ вҶ’ end_turn",
         "font": FONT_CODE, "size": 12, "color": MUTED,
         "space_after": 0},
    ])

    _add_text(s, 0.75, 6.6, 12, 0.4,
              "еҸҜиғҪиҰҒи·‘еҫҲеӨҡијӘ вҖ”вҖ” LLM еҸҜд»ҘйҖЈзәҢе‘јеҸ«еӨҡеҖӢе·Ҙе…·,зӣҙеҲ°е®ғиҰәеҫ—еӨ дәҶ (end_turn)гҖӮ",
              font=FONT_BODY, size=15, color=MUTED, italic=True,
              align=PP_ALIGN.CENTER)
    return s


def build_fc_takeaway(prs):
    """Takeaway: LLM вҶ” Harness responsibility split + transition to MCP."""
    s = _chrome_slide(prs)
    _add_text(s, 0.75, 0.4, 12, 0.7,
              "йҮҚй»һ:LLM вҶ” Harness иҒ·иІ¬еҲҶе·Ҙ",
              font=FONT_TITLE, size=28, color=NAVY, bold=True)
    _add_rect(s, 0.75, 1.20, 3, 0.04, ORANGE)

    # Left: LLM card
    llm_x, llm_w = 0.75, 5.9
    _add_rounded(s, llm_x, 1.7, llm_w, 4.5, WHITE,
                 line_color=BLUE, line_w=2.5)
    _add_rect(s, llm_x + 0.08, 1.78, llm_w - 0.16, 0.55, BLUE)
    _add_text(s, llm_x + 0.08, 1.78, llm_w - 0.16, 0.55, "LLM еҒҡзҡ„",
              font=FONT_TITLE, size=22, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_multi(s, llm_x + 0.4, 2.6, llm_w - 0.7, 3.4, [
        {"text": "вҖў зҙ”зІ№еҗҗеӯ—дёІ (text + tool_use JSON)",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вҖў и®Җ tool schema,з”ҹ input JSON",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вҖў Constrained decoding зўәдҝқ JSON еҗҲжі•",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вҖў дёҚзҹҘйҒ“е·Ҙе…·зңҹжӯЈеҒҡд»Җйәј",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вң— зў°дёҚеҲ° OSгҖҒжӘ”жЎҲгҖҒз¶Іи·Ҝ",
         "font": FONT_BODY, "size": 17, "color": RED, "bold": True,
         "space_after": 0},
    ])

    # Right: Harness card
    h_x, h_w = 6.95, 5.9
    _add_rounded(s, h_x, 1.7, h_w, 4.5, WHITE,
                 line_color=ORANGE, line_w=2.5)
    _add_rect(s, h_x + 0.08, 1.78, h_w - 0.16, 0.55, ORANGE)
    _add_text(s, h_x + 0.08, 1.78, h_w - 0.16, 0.55, "Harness еҒҡзҡ„",
              font=FONT_TITLE, size=22, color=WHITE, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_multi(s, h_x + 0.4, 2.6, h_w - 0.7, 3.4, [
        {"text": "вҖў зңҹзҡ„еҹ·иЎҢ (subprocessгҖҒHTTPгҖҒDBвҖҰ)",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вҖў з¶ӯиӯ· messages йҷЈеҲ—",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вҖў Loop зӣҙеҲ° stop_reason: end_turn",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вҖў иҷ•зҗҶ timeout / йҢҜиӘӨ / зөҗжһңж јејҸ",
         "font": FONT_BODY, "size": 17, "color": DARK, "space_after": 8},
        {"text": "вң“ еҠ  sandbox / е®үе…ЁйҳІиӯ·",
         "font": FONT_BODY, "size": 17, "color": GREEN, "bold": True,
         "space_after": 0},
    ])

    # Bottom transition to next section
    _add_text(s, 0.75, 6.5, 12, 0.5,
              "вҶ’ MCP еңЁйҖҷеҖӢж©ҹеҲ¶дёҠеҠ дәҶд»Җйәј?жҠҠ harness вҶ” tool зҡ„еҚ”е®ҡжЁҷжә–еҢ–гҖҒжҠҪеҲ°зҚЁз«Ӣ processгҖӮ",
              font=FONT_BODY, size=17, color=NAVY, bold=True, italic=True,
              align=PP_ALIGN.CENTER)
    return s


def _replace_exact_text(slide, old, new):
    """Replace any text run whose entire text == old, in place."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text == old:
                    run.text = new


def rebuild_architecture_slide(slide):
    """Replace the old HTML-demo cue with a static 4-layer diagram."""
    _clear_slide_shapes(slide)
    _apply_chrome(slide)
    _add_text(slide, 0.75, 0.4, 12, 0.7,
              "MCP ж•ҙй«”жһ¶ж§Ӣ:еӣӣеұӨеҲҶе·Ҙ",
              font=FONT_TITLE, size=28, color=NAVY, bold=True)
    _add_rect(slide, 0.75, 1.20, 3, 0.04, ORANGE)

    layers = [
        ("Host",   "жҮүз”ЁзЁӢејҸ / е…ҘеҸЈ",
         "иҲҲеӨ§ AI еӯёдјҙ web UI",                                  NAVY),
        ("Client", "еҢ…еҗ« LLM + MCP client йӮҸијҜ",
         "Node.js + Sonnet  (= еҪұзүҮиЈЎзҡ„ Parent)",                 BLUE),
        ("Server", "еҖӢеҲҘ MCP server еӯҗзЁӢеәҸ",
         "Python MCP Server  (= еҪұзүҮиЈЎзҡ„ Child)",                 CYAN),
        ("Tool",   "еҖӢеҲҘе·Ҙе…·еҮҪејҸ",
         "search_new_books / search_courses / search_teachers вҖҰ", GREEN),
    ]

    band_h = 1.00
    band_y = 1.7
    gap = 0.20
    for i, (name, desc, ex, color) in enumerate(layers):
        y = band_y + i * (band_h + gap)
        _add_rounded(slide, 0.8, y, 12, band_h, WHITE,
                     line_color=color, line_w=2.5)
        # Colored left strip
        _add_rect(slide, 0.85, y + 0.06, 0.2, band_h - 0.12, color)
        # Layer name
        _add_text(slide, 1.2, y, 2.5, band_h, name,
                  font=FONT_TITLE, size=24, color=color, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
        # Description (middle)
        _add_text(slide, 4.0, y, 4.3, band_h, desc,
                  font=FONT_BODY, size=18, color=DARK,
                  anchor=MSO_ANCHOR.MIDDLE)
        # Example (right)
        _add_text(slide, 8.5, y, 4.3, band_h, ex,
                  font=FONT_CODE, size=14, color=MUTED, italic=True,
                  anchor=MSO_ANCHOR.MIDDLE)

        # Arrow indicator below each band (except last)
        if i < len(layers) - 1:
            ay = y + band_h + 0.03
            _add_text(slide, 6.4, ay, 0.6, gap - 0.04, "в–ј",
                      font=FONT_BODY, size=16, color=MUTED,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text(slide, 0.75, 6.95, 12, 0.4,
              "Host зңӢдёҚеҲ° JSON-RPC В· Client зңӢдёҚеҲ° Server еҜҰдҪң В· LLM еҸӘзңӢ name + description + inputSchema",
              font=FONT_BODY, size=14, color=MUTED, italic=True,
              align=PP_ALIGN.CENTER)


def rebuild_recap(slide):
    """Rebuild final recap slide with refreshed vocab."""
    _clear_slide_shapes(slide)
    _apply_chrome(slide)
    _add_text(slide, 0.75, 0.5, 12, 0.8, "жң¬и¬ӣйҮҚй»һеӣһйЎ§",
              font=FONT_TITLE, size=36, color=NAVY, bold=True)
    _add_rect(slide, 0.75, 1.40, 3, 0.04, ORANGE)

    items = [
        "Function calling:LLM еҸӘеҗҗ tool_use JSON,зңҹзҡ„еҹ·иЎҢзҡ„жҳҜеӨ–йқўзҡ„ harness;harness и·Ёи¶ҠгҖҢж–Үеӯ—е®Үе®ҷ вҶ’ зңҹеҜҰйӣ»и…ҰгҖҚ",
        "MCP жҠҠ LLM вҮ„ е·Ҙе…·зҡ„йҖҡиЁҠжҠҪеҮәдҫҶ,и®ҠжҲҗ Parent / Child е…©еҖӢ process д№Ӣй–“зҡ„йӣҷеҗ‘з®ЎйҒ“",
        "е…©йҡҺж®өжҸЎжүӢ:в‘  дә’жҸӣиғҪеҠӣ вҶ’ в‘Ў жӢҝеҲ°е·Ҙе…·жё…е–®;е…©жӯҘйғҪйҒҺжүҚз®—зңҹзҡ„йҖЈдёҠ",
        "е·Ҙе…·е‘јеҸ«:LLM жұәе®ҡ вҶ’ token йЈӣеҺ» Child вҶ’ Child еҹ·иЎҢ вҶ’ зөҗжһңйЈӣеӣһ вҶ’ жүҫеҲ°е°ҚжҮүи«ӢжұӮеӣһзөҰдҪҝз”ЁиҖ…",
        "жҜҸеҖӢе·Ҙе…·з”ұ name / description / inputSchema дёүеҖӢж¬„дҪҚе®ҡзҫ©;description зҡ„е“ҒиіӘзӣҙжҺҘжұәе®ҡ LLM йҒёдёҚйҒёеҫ—е°Қ",
        "Client жҠҠжүҖжңүе·Ҙе…·е®ҡзҫ©иҪүжҸӣеҫҢ,жіЁе…Ҙ LLM API зҡ„ tools еҸғж•ё вҖ” LLM еҸӘзңӢеҫ—еҲ° name + description + inputSchema",
    ]
    y = 1.65
    for i, txt in enumerate(items, 1):
        _add_text(slide, 0.9, y, 0.8, 0.5, str(i),
                  font=FONT_TITLE, size=22, color=ORANGE, bold=True,
                  anchor=MSO_ANCHOR.MIDDLE)
        _add_text(slide, 1.8, y, 11.3, 0.5, txt,
                  font=FONT_BODY, size=16, color=DARK,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += 0.78

    _add_text(slide, 0.75, 6.8, 12, 0.4,
              "дёӢдёҖи¬ӣй җе‘Ҡ:Agentic Tool Loop вҖ” LLM жҖҺйәјиҮӘдё»жұәе®ҡи©Іе‘јеҸ«е“ӘеҖӢе·Ҙе…·",
              font=FONT_BODY, size=16, color=TEAL, italic=True)


# в”Җв”Җ Main flow в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def main():
    prs = Presentation(str(PPTX))
    n_orig = len(prs.slides)
    print(f"opened {PPTX.name}: {n_orig} slides")

    # Build new slides вҖ” appended at end first
    new_indices = {}

    # Section 01: Function Calling жҖҺйәјйҒӢдҪң (new)
    new_indices["fc_divider"]  = len(prs.slides); build_slide_section_divider(
        prs, "01", "Function Calling жҖҺйәјйҒӢдҪң",
        "LLM еҸӘеҗҗеӯ—дёІ,зңҹзҡ„еҹ·иЎҢзҡ„жҳҜеӨ–йқўзҡ„ harness")
    new_indices["fc_setup"]    = len(prs.slides); build_fc_setup(prs)
    new_indices["fc_stage12"]  = len(prs.slides); build_fc_stage12(prs)
    new_indices["fc_stage3"]   = len(prs.slides); build_fc_stage3(prs)
    new_indices["fc_stage4"]   = len(prs.slides); build_fc_stage4(prs)
    new_indices["fc_stage56"]  = len(prs.slides); build_fc_stage56(prs)
    new_indices["fc_takeaway"] = len(prs.slides); build_fc_takeaway(prs)

    # Section 02: еҫһдёҖеҖӢжҹҘи©ўй–Ӣе§Ӣ (renumbered from previous 01)
    new_indices["divider"]  = len(prs.slides); build_slide_section_divider(
        prs, "02", "еҫһдёҖеҖӢжҹҘи©ўй–Ӣе§Ӣ",
        "з”Ё search_new_books иө°дёҖйҒҚе®Ңж•ҙжөҒзЁӢ")
    new_indices["scenario"] = len(prs.slides); build_slide_scenario(prs)
    new_indices["act1"]     = len(prs.slides); build_slide_act1(prs)
    new_indices["act2"]     = len(prs.slides); build_slide_act2(prs)
    new_indices["act3"]     = len(prs.slides); build_slide_act3(prs)
    new_indices["video"]    = len(prs.slides); build_slide_video_cue(prs)

    # Appendix divider
    new_indices["appendix"] = len(prs.slides); build_appendix_divider(prs)

    print(f"appended 14 new slides; total now {len(prs.slides)}")

    # Reorder. Goal final order:
    #   1: orig slide 1 (title)
    #   2: orig slide 2 (outline, REBUILT)
    #   3: new section divider
    #   4: new scenario
    #   5: new act1
    #   6: new act2
    #   7: new act3
    #   8: new video cue
    #   9-12: orig slide 12-15 (Tool description)
    #   13-16: orig slide 16-19 (Client integration)
    #   17: orig slide 20 (recap, REBUILT)
    #   18: new appendix divider
    #   19-23: orig slide 3-7 (REST vs JSON-RPC content)
    # DELETE: orig slide 8-11 (lifecycle)

    # Use sldId-element references to be safe across reorders
    xml = _slides_xml(prs)
    all_slides = list(xml)
    orig_3_to_7   = all_slides[2:7]
    orig_8_to_11  = all_slides[7:11]
    orig_12_to_19 = all_slides[11:19]
    orig_20       = all_slides[19]

    # New FC section (7 slides, indices 20..26 in the appended order)
    new_fc_divider  = all_slides[20]
    new_fc_setup    = all_slides[21]
    new_fc_stage12  = all_slides[22]
    new_fc_stage3   = all_slides[23]
    new_fc_stage4   = all_slides[24]
    new_fc_stage56  = all_slides[25]
    new_fc_takeaway = all_slides[26]

    # New query section (6 slides, indices 27..32)
    new_divider  = all_slides[27]
    new_scenario = all_slides[28]
    new_act1     = all_slides[29]
    new_act2     = all_slides[30]
    new_act3     = all_slides[31]
    new_video    = all_slides[32]

    # Appendix divider (index 33)
    new_appendix = all_slides[33]

    # Remove all from XML; then re-append in desired order
    for el in list(xml):
        xml.remove(el)

    final = [
        all_slides[0],     # 1: title (kept)
        all_slides[1],     # 2: outline (REBUILT)

        # Section 01 вҖ” Function Calling жҖҺйәјйҒӢдҪң (new)
        new_fc_divider,
        new_fc_setup,
        new_fc_stage12,
        new_fc_stage3,
        new_fc_stage4,
        new_fc_stage56,
        new_fc_takeaway,

        # Section 02 вҖ” еҫһдёҖеҖӢжҹҘи©ўй–Ӣе§Ӣ (renumbered)
        new_divider,
        new_scenario,
        new_act1,
        new_act2,
        new_act3,
        new_video,

        # Sections 03 + 04 вҖ” Tool иЁ»еҶҠиҲҮжҸҸиҝ° / Client ж•ҙеҗҲж©ҹеҲ¶ (kept original)
        *orig_12_to_19,

        # Recap
        orig_20,           # REBUILT

        # Appendix
        new_appendix,
        *orig_3_to_7,
    ]
    # orig_8_to_11 (lifecycle) dropped

    for el in final:
        xml.append(el)
    print(f"reordered to {len(final)} slides (dropped 4 lifecycle slides)")

    # NOTE on deletion: removing the sldId reference unhooks the slide from
    # the deck but the underlying slide part still lives in the package. Good
    # enough for our purposes вҖ” file is slightly larger but no broken refs.

    # Rebuild outline slide (slide 2 вҶ’ index 1)
    rebuild_slide_2(prs.slides[1])
    # Recap is now at index 23 (24th slide, 0-based) after the FC + query
    # sections + 8 kept slides above it
    rebuild_recap(prs.slides[23])
    print("rebuilt outline (slide 2) and recap (slide 24)")

    # Kept section dividers already have correct numbers (03 Tool / 04 Client)
    # from the original pptx вҖ” no renumbering needed in this revision.

    # Replace the old "вҶ’ еҲҮжҸӣеҲ° mcp-architecture-animation.html" demo cue
    # (originally orig slide 19, now at index 22) with a static 4-layer diagram.
    rebuild_architecture_slide(prs.slides[22])
    print("rebuilt architecture slide (slide 23) вҖ” no more HTML demo cue")

    prs.save(str(PPTX))
    print(f"saved вҶ’ {PPTX.name} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
